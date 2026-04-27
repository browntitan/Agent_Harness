"""
Apply audit updates to the PowerPoint deck:
1. Global text corrections (endpoint counts, agent tool counts, etc.)
2. Add new slides for missing subsystems
3. Update speaker notes
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn
import copy, os

# ─── CONSTANTS ───────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
TEAL      = RGBColor(0x00, 0x97, 0xA7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEAL = RGBColor(0x00, 0x6D, 0x77)
GOLD      = RGBColor(0xFF, 0xB7, 0x4D)
MED_GRAY  = RGBColor(0x99, 0x99, 0x99)
BG_LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)
CARD_BG   = RGBColor(0x23, 0x3B, 0x5E)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
ACCENT_PURPLE = RGBColor(0x7C, 0x4D, 0xFF)
ACCENT_RED = RGBColor(0xEF, 0x53, 0x50)
ACCENT_BLUE = RGBColor(0x42, 0xA5, 0xF5)

SW = Inches(10)
SH = Inches(5.625)
MARGIN = Inches(0.55)

# ─── HELPER FUNCTIONS ───────────────────────────────────────
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                font_name="Arial", valign=MSO_ANCHOR.TOP, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    tf.auto_size = None
    txBox.text_frame._txBody.bodyPr.set("anchor", {
        MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"
    }.get(valign, "t"))
    if line_spacing:
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txBox

def add_rich_textbox(slide, left, top, width, height, runs, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    current_p = tf.paragraphs[0]
    first_run = True
    for run_def in runs:
        if run_def.get('breakLine') and not first_run:
            current_p = tf.add_paragraph()
        if run_def.get('align'):
            current_p.alignment = run_def['align']
        r = current_p.add_run()
        r.text = run_def.get('text', '')
        r.font.size = Pt(run_def.get('font_size', 14))
        r.font.bold = run_def.get('bold', False)
        r.font.color.rgb = run_def.get('color', WHITE)
        r.font.name = run_def.get('font_name', 'Arial')
        if line_spacing:
            pPr = current_p._pPr if current_p._pPr is not None else current_p._p.get_or_add_pPr()
            if pPr.find(qn('a:lnSpc')) is None:
                lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
                spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
                spcPts.set('val', str(int(line_spacing * 100)))
        first_run = False
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body, accent_color=TEAL,
             title_size=16, body_size=12, bg_color=CARD_BG):
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill_color=bg_color)
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height, fill_color=accent_color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.3), Inches(0.4),
                title, font_size=title_size, bold=True, color=WHITE)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6),
                body, font_size=body_size, color=LIGHT_GRAY, line_spacing=18)

def build_section_divider(slide, title, subtitle):
    set_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, Inches(0.06), fill_color=TEAL)
    add_textbox(slide, MARGIN, Inches(1.8), Inches(9), Inches(1.0),
                title, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(slide, MARGIN, Inches(2.8), Inches(8), Inches(0.8),
                subtitle, font_size=16, color=TEAL, align=PP_ALIGN.LEFT)

def build_content_slide(slide, title, body_callback):
    set_bg(slide, BG_LIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, Inches(0.75), fill_color=NAVY)
    add_textbox(slide, MARGIN, Inches(0.12), Inches(9), Inches(0.5),
                title, font_size=22, bold=True, color=WHITE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.75), SW, Inches(0.04), fill_color=TEAL)
    body_callback(slide)

def add_notes(slide, text):
    """Add speaker notes to a slide."""
    from pptx.oxml.ns import qn as _qn
    # Ensure notes slide exists
    if not slide.has_notes_slide:
        # Access the notes slide to create it
        _ = slide.notes_slide
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    if tf is not None:
        tf.text = text
    else:
        # Fallback: create notes body manually
        notes_body = notes_slide.notes_placeholder
        if notes_body is not None:
            notes_body.text = text

def set_notes(slide, text):
    """Set speaker notes, replacing existing."""
    if not slide.has_notes_slide:
        _ = slide.notes_slide
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    if tf is not None:
        for p in tf.paragraphs:
            p.text = ""
        tf.paragraphs[0].text = text
    else:
        notes_body = notes_slide.notes_placeholder
        if notes_body is not None:
            notes_body.text = text


# ══════════════════════════════════════════════════════════════
# STEP 1: GLOBAL TEXT REPLACEMENTS
# ══════════════════════════════════════════════════════════════

SRC = os.path.join(os.path.dirname(__file__), "agentic_chatbot_v2_deep_dive_updated.pptx")
OUT = os.path.join(os.path.dirname(__file__), "agentic_chatbot_v3_audit_updated.pptx")

prs = Presentation(SRC)

# Replacement map for text corrections
REPLACEMENTS = {
    # API endpoint count corrections
    "47 public endpoints": "48 public endpoints",
    "47 Public Endpoints": "48 Public Endpoints",
    "PUBLIC API (47 endpoints)": "PUBLIC API (48 endpoints)",
    "117 Endpoints": "118 Endpoints",
    "117 total API surface": "118 total API surface",
    "api/main.py \u2014 47 public endpoints": "api/main.py \u2014 48 public endpoints",
    # General agent description (exact context)
    "general.md (19 tools)": "general.md (33 tools)",
    "coordinator.md (4 tools)": "coordinator.md (11 tools)",
    "data_analyst.md (14 tools)": "data_analyst.md (20 tools)",
    # Coordinator allowed workers update
    "8 allowed workers": "9 allowed workers",
    "Allowed workers: general, utility, data_analyst, rag_worker, planner, finalizer, verifier, memory_maintainer":
        "Allowed workers: general, utility, data_analyst, rag_worker, graph_manager, planner, finalizer, verifier, memory_maintainer",
}

# Additional targeted replacements for notes
NOTES_REPLACEMENTS = {
    "API has grown to 23 endpoints": "API has grown to 48 endpoints",
    "23 Public Endpoints": "48 Public Endpoints",
}

def apply_text_replacements(prs, replacements):
    """Apply text replacements across all slides, shapes, and notes."""
    count = 0
    for slide in prs.slides:
        # Replace in shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for old, new in replacements.items():
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                count += 1
        # Replace in notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            for paragraph in notes_tf.paragraphs:
                for run in paragraph.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                            count += 1
    return count

def apply_targeted_run_replacements(prs):
    """Apply precise per-run text replacements for agent tool counts.
    These need exact matching to avoid corrupting '54 tools' or '14 tools'."""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        t = run.text
                        # Fix "19 tools" only when it's the entire run text (general agent)
                        if t.strip() == "19 tools":
                            run.text = t.replace("19 tools", "33 tools")
                            count += 1
                        # Fix "4 tools" only when it's the entire run text (coordinator)
                        elif t.strip() == "4 tools":
                            run.text = t.replace("4 tools", "11 tools")
                            count += 1
                        # Fix "7 tools" only when it's the entire run text (utility)
                        elif t.strip() == "7 tools":
                            run.text = t.replace("7 tools", "13 tools")
                            count += 1
                        # Fix "5 tools" only when it's the entire run text (graph_manager)
                        elif t.strip() == "5 tools":
                            run.text = t.replace("5 tools", "12 tools")
                            count += 1
                        # Fix "Tools: 11" -> "Tools: 33" (general deep dive slide)
                        elif t.strip() == "Tools: 11":
                            run.text = t.replace("Tools: 11", "Tools: 33")
                            count += 1
                        # Fix "Tools: 4 (orchestration)" (coordinator deep dive)
                        elif "Tools: 4" in t and "orchestration" in t:
                            run.text = t.replace("Tools: 4", "Tools: 11")
                            count += 1
                        # Fix "Tools: 7" (utility)
                        elif t.strip() == "Tools: 7":
                            run.text = t.replace("Tools: 7", "Tools: 13")
                            count += 1
                        # Fix "Tools: 10" (verifier)
                        elif t.strip() == "Tools: 10":
                            run.text = t.replace("Tools: 10", "Tools: 13")
                            count += 1
                        # Fix "Tools: 5" (graph_manager)
                        elif t.strip() == "Tools: 5":
                            run.text = t.replace("Tools: 5", "Tools: 12")
                            count += 1
                        # Fix "10 tools" only when standalone (verifier on overview slide)
                        elif t.strip() == "10 tools":
                            run.text = t.replace("10 tools", "13 tools")
                            count += 1
                        # Fix "14 tools" only when standalone (data_analyst on overview slide)
                        elif t.strip() == "14 tools":
                            run.text = t.replace("14 tools", "20 tools")
                            count += 1
                        # Fix "Tools: 13" followed by Fact checker context (data analyst deep dive)
                        elif t.strip() == "Tools: 13":
                            run.text = t.replace("Tools: 13", "Tools: 20")
                            count += 1
        # Also fix in notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                for paragraph in notes_tf.paragraphs:
                    for run in paragraph.runs:
                        t = run.text
                        if "14 tools, 4 dedicated" in t:
                            run.text = t.replace("14 tools", "20 tools")
                            count += 1
    return count

print("Applying global text replacements...")
n = apply_text_replacements(prs, REPLACEMENTS)
print(f"  {n} replacements made in slides")
n2 = apply_text_replacements(prs, NOTES_REPLACEMENTS)
print(f"  {n2} replacements made in notes")
n3 = apply_targeted_run_replacements(prs)
print(f"  {n3} targeted tool count replacements")


# ══════════════════════════════════════════════════════════════
# STEP 2: ADD NEW SLIDES FOR MISSING SUBSYSTEMS
# ══════════════════════════════════════════════════════════════

layout = prs.slide_layouts[0]
new_slides = []

# ─── NEW SLIDE: Authorization Subsystem ──────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_authz(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.5),
                "Fine-grained access control with principals, roles, bindings, and permissions",
                font_size=13, color=RGBColor(0x66, 0x66, 0x66))

    # Left column - RBAC model
    add_textbox(s, MARGIN, Inches(1.5), Inches(4.2), Inches(0.35),
                "RBAC MODEL", font_size=13, bold=True, color=TEAL)

    items_left = [
        ("Principals", "Users and service accounts with unique IDs"),
        ("Roles", "Named permission bundles (admin, editor, viewer)"),
        ("Bindings", "Map principals to roles on specific resources"),
        ("Memberships", "Group-level role inheritance"),
        ("Permissions", "Granular resource:action pairs"),
    ]
    y = Inches(1.9)
    for title, desc in items_left:
        add_shape(s, MSO_SHAPE.OVAL, MARGIN, y + Inches(0.05), Inches(0.14), Inches(0.14), fill_color=TEAL)
        add_textbox(s, MARGIN + Inches(0.22), y, Inches(1.4), Inches(0.28),
                    title, font_size=11, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, Inches(2.2), y, Inches(2.5), Inches(0.28),
                    desc, font_size=10, color=RGBColor(0x66, 0x66, 0x66))
        y += Inches(0.38)

    # Right column - API surface
    add_textbox(s, Inches(5.3), Inches(1.5), Inches(4.2), Inches(0.35),
                "10 ADMIN ENDPOINTS", font_size=13, bold=True, color=ACCENT_ORANGE)

    endpoints = [
        "GET/POST  /admin/access/principals",
        "GET/POST  /admin/access/roles",
        "GET/POST  /admin/access/bindings",
        "GET/POST  /admin/access/memberships",
        "GET       /admin/access/permissions",
        "GET       /admin/access/effective-access",
    ]
    y = Inches(1.95)
    for ep in endpoints:
        add_textbox(s, Inches(5.3), y, Inches(4.2), Inches(0.25),
                    ep, font_size=10, color=RGBColor(0x44, 0x44, 0x44), font_name="Consolas")
        y += Inches(0.3)

    # Bottom box
    add_shape(s, MSO_SHAPE.RECTANGLE, MARGIN, Inches(4.2), Inches(8.9), Inches(1.1), fill_color=NAVY)
    add_textbox(s, Inches(0.8), Inches(4.35), Inches(8.4), Inches(0.8),
                "authz/service.py provides tenant-scoped resource access. "
                "Integrated with API gateway (gateway_security.py) and connector endpoints (connector_security.py). "
                "Per-skill access filtering in SkillRuntime. Tool policy enforcement in tools/policy.py.",
                font_size=11, color=LIGHT_GRAY, line_spacing=18)

build_content_slide(slide, "Authorization Subsystem (authz/)", body_authz)
add_notes(slide,
    "The authorization subsystem is a dedicated directory at authz/ providing fine-grained access control. "
    "It implements a standard RBAC model with principals, roles, bindings, memberships, and permissions. "
    "There are 10 admin endpoints for managing access control. "
    "It integrates with the API gateway for request-level auth, connector security for third-party integrations, "
    "SkillRuntime for per-skill access filtering, and the tool policy layer. "
    "This is separate from the API-level security and provides resource-level granularity.")


# ─── NEW SLIDE: Documents Subsystem Deep Dive ────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_docs(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "13 files across extractors, compare, consolidation, similarity, evidence, templates, and serializers",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # Three columns for the three main services
    col_w = Inches(2.8)
    gap = Inches(0.15)

    # Column 1: Extraction
    x1 = MARGIN
    add_card(s, x1, Inches(1.5), col_w, Inches(1.7),
             "Document Extraction",
             "9 formats: PDF, DOCX, PPTX, XLSX, XLS, TXT, MD, CSV, TSV. "
             "Docling parser for PDFs. Max 200 elements per extraction.",
             accent_color=TEAL, title_size=13, body_size=10)

    # Column 2: Comparison
    x2 = x1 + col_w + gap
    add_card(s, x2, Inches(1.5), col_w, Inches(1.7),
             "Document Comparison",
             "Key-based + fuzzy matching (0.72 threshold). "
             "14 obligation modalities. Binding strength ranking. "
             "Severity escalation for changed obligations.",
             accent_color=ACCENT_ORANGE, title_size=13, body_size=10)

    # Column 3: Consolidation
    x3 = x2 + col_w + gap
    add_card(s, x3, Inches(1.5), col_w, Inches(1.7),
             "Consolidation Campaigns",
             "Multi-doc corpus analysis with IDF weighting. "
             "6 focus modes. Union-find clustering. "
             "Weighted Jaccard scoring.",
             accent_color=ACCENT_PURPLE, title_size=13, body_size=10)

    # Bottom row: Evidence + Similarity
    add_card(s, MARGIN, Inches(3.4), Inches(4.35), Inches(1.2),
             "Evidence Binder",
             "Compile graded evidence into formatted reference packets with inline citations.",
             accent_color=ACCENT_GREEN, title_size=13, body_size=10)

    add_card(s, Inches(5.15), Inches(3.4), Inches(4.35), Inches(1.2),
             "Similarity Analysis",
             "TF-IDF fingerprinting. Process keyword matching (28 patterns). N-gram comparison.",
             accent_color=DARK_TEAL, title_size=13, body_size=10)

build_content_slide(slide, "Documents Subsystem Deep Dive (documents/)", body_docs)
add_notes(slide,
    "The documents subsystem is far more sophisticated than simple extraction. "
    "DocumentExtractionService supports 9 file formats with Docling as the primary PDF parser. "
    "Each extraction produces structured output: sections with IDs (sec_XXXX), elements (el_XXXXX), tables, and figures. "
    "DocumentComparisonService does clause-level diffing with 14 obligation modality patterns for legal/compliance use cases. "
    "It tracks binding strength changes: if a 'may' becomes a 'shall', that's flagged as a severity escalation. "
    "DocumentConsolidationCampaignService runs corpus-wide analysis to find duplicate policies, harmonize processes, "
    "and identify consolidation opportunities. Uses union-find clustering and weighted Jaccard with context-specific scoring. "
    "The Evidence Binder compiles graded evidence into reference packets. "
    "Similarity analysis uses TF-IDF fingerprinting with n-grams and process keyword detection.")


# ─── NEW SLIDE: Skills Architecture Deep Dive ────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_skills_deep(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "10 files: pack_loader, resolver, runtime, indexer, telemetry, dependency_graph, execution, query_builder",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # Top row: 3 key concepts
    concepts = [
        ("3 Skill Kinds", "retrievable\nexecutable\nhybrid", TEAL),
        ("2 Exec Contexts", "inline (default)\nfork (isolated)", ACCENT_ORANGE),
        ("5 Effort Levels", "(empty), low\nmedium, high, xhigh", ACCENT_PURPLE),
    ]
    x = MARGIN
    for title, body, color in concepts:
        add_shape(s, MSO_SHAPE.RECTANGLE, x, Inches(1.4), Inches(2.8), Inches(1.1), fill_color=color)
        add_textbox(s, x + Inches(0.15), Inches(1.48), Inches(2.5), Inches(0.3),
                    title, font_size=14, bold=True, color=WHITE)
        add_textbox(s, x + Inches(0.15), Inches(1.8), Inches(2.5), Inches(0.6),
                    body, font_size=11, color=RGBColor(0xEE, 0xEE, 0xEE))
        x += Inches(3.05)

    # Bottom: Architecture flow
    add_textbox(s, MARGIN, Inches(2.7), Inches(4), Inches(0.35),
                "SKILL LIFECYCLE", font_size=13, bold=True, color=NAVY)

    flow_items = [
        ("SkillPackFile", "YAML frontmatter + markdown body. SHA256 checksum. Tags: tool_tags, task_tags."),
        ("SkillIndexSync", "Discovers .md files, chunks for RAG, builds inverted index. Incremental sync by checksum."),
        ("SkillResolver", "Semantic retrieval with BM25/embedding. Filters by agent_scope, tags. Pinned skill priority."),
        ("SkillRuntime", "Merges query + skills + session context. Access control filtering. Max 4000 chars budget."),
    ]
    y = Inches(3.1)
    for i, (title, desc) in enumerate(flow_items):
        num_color = TEAL if i % 2 == 0 else ACCENT_ORANGE
        add_shape(s, MSO_SHAPE.OVAL, MARGIN, y + Inches(0.03), Inches(0.22), Inches(0.22), fill_color=num_color)
        add_textbox(s, MARGIN + Inches(0.02), y + Inches(0.03), Inches(0.22), Inches(0.22),
                    str(i+1), font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, MARGIN + Inches(0.3), y, Inches(1.6), Inches(0.25),
                    title, font_size=11, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, Inches(2.5), y, Inches(7), Inches(0.25),
                    desc, font_size=10, color=RGBColor(0x55, 0x55, 0x55))
        y += Inches(0.35)

    # Right side: Telemetry + Dependency
    add_card(s, Inches(5.3), Inches(4.55), Inches(4.2), Inches(0.85),
             "Telemetry + Dependencies",
             "80% success SLO, 20-use review window. DFS cycle detection. States: healthy/unstable/broken.",
             accent_color=ACCENT_GREEN, title_size=12, body_size=10)

    add_card(s, MARGIN, Inches(4.55), Inches(4.5), Inches(0.85),
             "Legacy Aliases + Visibility",
             "9 tool tag aliases mapped. Visibility: global / private / tenant scoped.",
             accent_color=MED_GRAY, title_size=12, body_size=10)

build_content_slide(slide, "Skills Architecture Deep Dive (skills/)", body_skills_deep)
add_notes(slide,
    "The skills subsystem has 10 files and is much more sophisticated than the previous slides showed. "
    "Three skill kinds: retrievable (injected as guidance), executable (run with specific tools), hybrid (both). "
    "Two execution contexts: inline runs within the current agent turn, fork runs in isolation. "
    "Five effort levels control how much compute a skill gets. "
    "SkillPackFile parses YAML frontmatter with SHA256 checksums. Fields include tool_tags, task_tags, when_to_apply, avoid_when, examples. "
    "SkillIndexSync discovers markdown files and builds an inverted index, with incremental sync by checksum so only changed skills re-index. "
    "SkillResolver does semantic retrieval with BM25 and embedding scoring, filtered by agent_scope and tags. Pinned skills get priority. "
    "SkillRuntime applies access control and merges everything into a 4000-char budget. "
    "SkillTelemetry tracks answer quality with an 80% success SLO over a 20-use review window. "
    "SkillDependencyGraph validates skill dependencies using DFS cycle detection with three health states.")


# ─── NEW SLIDE: Provider System & Circuit Breaker ────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_providers(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "6 files: factory, llm_factory, circuit_breaker, output_limits, dependency_checks",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # Left: Provider Resolution
    add_textbox(s, MARGIN, Inches(1.5), Inches(4.3), Inches(0.35),
                "PROVIDER RESOLUTION", font_size=13, bold=True, color=TEAL)

    add_card(s, MARGIN, Inches(1.9), Inches(4.3), Inches(1.4),
             "AgentProviderResolver",
             "Caches bundles by 6-tuple (provider, model, cap). "
             "Per-agent overrides: AGENT_{NAME}_CHAT_MODEL. "
             "Token cap chain: request > agent > global > ollama.",
             accent_color=TEAL, title_size=12, body_size=10)

    add_card(s, MARGIN, Inches(3.45), Inches(4.3), Inches(0.95),
             "ProviderBundle",
             "Chat + judge + embeddings. Dims: 768 (Nomic) / 1536 (Ada-002). "
             "Timeouts: 120s HTTP, 20s connect.",
             accent_color=DARK_TEAL, title_size=12, body_size=10)

    # Right: Circuit Breaker
    add_textbox(s, Inches(5.3), Inches(1.5), Inches(4.2), Inches(0.35),
                "CIRCUIT BREAKER", font_size=13, bold=True, color=ACCENT_RED)

    add_shape(s, MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(1.9), Inches(4.2), Inches(2.45), fill_color=CARD_BG)
    add_shape(s, MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(1.9), Inches(0.06), Inches(2.45), fill_color=ACCENT_RED)

    breaker_items = [
        ("Window Size:", "20 samples"),
        ("Min Samples:", "6 before tripping"),
        ("Error Threshold:", "50% error rate"),
        ("Consecutive Fails:", "3 triggers open"),
        ("Open Duration:", "30 seconds cooldown"),
        ("Fallback:", "Degrades to basic mode"),
    ]
    y = Inches(2.05)
    for label, val in breaker_items:
        add_textbox(s, Inches(5.55), y, Inches(1.8), Inches(0.25),
                    label, font_size=11, bold=True, color=WHITE)
        add_textbox(s, Inches(7.4), y, Inches(2.0), Inches(0.25),
                    val, font_size=11, color=LIGHT_GRAY)
        y += Inches(0.35)

    # Bottom
    add_shape(s, MSO_SHAPE.RECTANGLE, MARGIN, Inches(4.6), Inches(8.9), Inches(0.8), fill_color=NAVY)
    add_textbox(s, Inches(0.8), Inches(4.7), Inches(8.4), Inches(0.6),
                "If an agent's LLM provider trips the circuit breaker, RuntimeService degrades to basic mode "
                "with a different provider. If even that fails, returns a graceful degraded message to the user.",
                font_size=11, color=LIGHT_GRAY, line_spacing=18)

build_content_slide(slide, "Provider System & Circuit Breaker (providers/)", body_providers)
add_notes(slide,
    "The providers directory has 6 files managing LLM provider resolution and fault tolerance. "
    "AgentProviderResolver caches provider bundles keyed by a 6-tuple of provider/model/cap settings. "
    "Per-agent model overrides let you run different agents on different models via environment variables like AGENT_DATA_ANALYST_CHAT_MODEL. "
    "Output token capping follows a priority chain: explicit request cap, then agent-specific, then global/demo, then ollama_num_predict. "
    "The circuit breaker is critical for production reliability. It monitors a sliding window of 20 samples. "
    "If the error rate exceeds 50% or 3 consecutive failures occur, the circuit opens for 30 seconds. "
    "During that time, RuntimeService falls back to basic mode with a different provider. "
    "This prevents cascading failures when an LLM provider has an outage.")


# ─── NEW SLIDE: MCP Security Model ──────────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_mcp_security(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "Cryptographic isolation, URL validation, tool deferral, and connection lifecycle management",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # Four cards in 2x2 grid
    cw = Inches(4.35)
    ch = Inches(1.5)

    add_card(s, MARGIN, Inches(1.5), cw, ch,
             "Secret Encryption",
             "Fernet symmetric encryption with SHA256 key derivation. "
             "Prefix: 'fernet:v1:' identifies encrypted values. "
             "Secrets never stored in plaintext. Decrypted only at call time.",
             accent_color=ACCENT_RED, title_size=13, body_size=10)

    add_card(s, Inches(5.15), Inches(1.5), cw, ch,
             "URL Validation",
             "HTTPS enforced by default (configurable). "
             "Private network blocking prevents SSRF attacks. "
             "Connection names slugified: max 80 chars, [a-z0-9_] only.",
             accent_color=ACCENT_ORANGE, title_size=13, body_size=10)

    add_card(s, MARGIN, Inches(3.2), cw, ch,
             "Tool Deferral System",
             "All MCP tools: defer=True, destructive=True, defer_priority=50. "
             "Tools only execute after explicit user/agent approval. "
             "Schema sanitization: enforces type='object', additionalProperties=True.",
             accent_color=TEAL, title_size=13, body_size=10)

    add_card(s, Inches(5.15), Inches(3.2), cw, ch,
             "Connection Lifecycle",
             "Timeouts: 15s for listing tools, 60s for execution. "
             "Tool name pattern: mcp__{connection}__{tool}. "
             "Per-tool enable/disable. Status/enabled flags gate execution.",
             accent_color=ACCENT_PURPLE, title_size=13, body_size=10)

    # Bottom note
    add_textbox(s, MARGIN, Inches(4.9), Inches(8.9), Inches(0.5),
                "McpStreamableHttpClient wraps the official MCP Streamable HTTP client with sync/async support. "
                "McpCatalogService manages the full connection CRUD + tool discovery + execution pipeline.",
                font_size=10, color=RGBColor(0x88, 0x88, 0x88))

build_content_slide(slide, "MCP Security Model (mcp/)", body_mcp_security)
add_notes(slide,
    "The MCP subsystem has 4 files implementing a secure Model Context Protocol integration. "
    "All MCP secrets are encrypted with Fernet symmetric encryption using SHA256-derived keys. "
    "The 'fernet:v1:' prefix identifies encrypted values. Secrets are only decrypted at tool call time. "
    "URL validation enforces HTTPS and blocks private networks to prevent SSRF attacks. "
    "The tool deferral system marks all MCP tools as deferred and destructive by default, "
    "meaning they require explicit approval before execution. Schema sanitization ensures consistent tool interfaces. "
    "Connection timeouts are split: 15 seconds for listing available tools, 60 seconds for actual execution. "
    "Tool names follow the mcp__{connection}__{tool} pattern for namespace isolation.")


# ─── NEW SLIDE: Task Plan & Artifact Handoff ─────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_task_plan(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "Structured task decomposition with typed artifact handoffs between worker agents",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # Left: TaskSpec
    add_textbox(s, MARGIN, Inches(1.4), Inches(4.3), Inches(0.35),
                "TASK SPECIFICATION", font_size=13, bold=True, color=TEAL)

    spec_items = [
        ("Executors:", "rag_worker, utility, data_analyst, general, graph_manager, verifier"),
        ("Modes:", "sequential or parallel execution"),
        ("Terminal States:", "completed, failed, stopped, waiting_message"),
        ("Hints:", "research_profile, coverage_goal, result_mode, answer_mode"),
    ]
    y = Inches(1.85)
    for label, val in spec_items:
        add_textbox(s, MARGIN, y, Inches(1.5), Inches(0.25),
                    label, font_size=11, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, Inches(2.0), y, Inches(2.8), Inches(0.25),
                    val, font_size=10, color=RGBColor(0x55, 0x55, 0x55))
        y += Inches(0.35)

    # Right: Artifact Types
    add_textbox(s, Inches(5.3), Inches(1.4), Inches(4.2), Inches(0.35),
                "8 ARTIFACT HANDOFF TYPES", font_size=13, bold=True, color=ACCENT_ORANGE)

    artifacts = [
        "title_candidates",
        "doc_focus",
        "research_facets",
        "facet_matches",
        "doc_digest",
        "subsystem_inventory",
        "policy_guidance_matches",
        "buyer_recommendation_table",
    ]
    y = Inches(1.85)
    for art in artifacts:
        add_shape(s, MSO_SHAPE.RECTANGLE, Inches(5.3), y, Inches(0.12), Inches(0.18), fill_color=ACCENT_ORANGE)
        add_textbox(s, Inches(5.55), y, Inches(3.9), Inches(0.22),
                    art, font_size=11, color=RGBColor(0x33, 0x33, 0x33), font_name="Consolas")
        y += Inches(0.28)

    # Bottom: Document Ranking
    add_shape(s, MSO_SHAPE.RECTANGLE, MARGIN, Inches(4.2), Inches(8.9), Inches(1.1), fill_color=NAVY)
    add_textbox(s, Inches(0.8), Inches(4.3), Inches(8.4), Inches(0.35),
                "DOCUMENT RANKING KEYS", font_size=12, bold=True, color=TEAL)
    add_textbox(s, Inches(0.8), Inches(4.65), Inches(8.4), Inches(0.5),
                "is_meta_document  |  reviewed_relevance  |  matched_facets  |  strong_evidence_count  |  title_path_score  |  seed_hits",
                font_size=11, color=LIGHT_GRAY, font_name="Consolas")

build_content_slide(slide, "Task Plan & Artifact Handoff System", body_task_plan)
add_notes(slide,
    "The task plan system in runtime/task_plan.py and kernel_coordinator.py manages structured multi-agent workflows. "
    "TaskSpec defines what a task does: which executor runs it, whether it's sequential or parallel, "
    "and controller hints like research_profile, coverage_goal, and result_mode. "
    "There are 6 valid executor types: rag_worker, utility, data_analyst, general, graph_manager, and verifier. "
    "The artifact handoff system enables typed data passing between workers. "
    "8 artifact types cover the full range from title candidates to buyer recommendation tables. "
    "Each artifact type has an allowlist of valid consumers, preventing incorrect data flow. "
    "KernelCoordinatorController builds ranked document lists from artifacts using 6 scoring dimensions: "
    "is_meta_document, reviewed_relevance, matched_facets, strong_evidence_count, title_path_score, and seed_hits. "
    "A meta-document filter uses regex to exclude test fixtures and acceptance scenarios from ranking.")


# ─── NEW SLIDE: Storage & Blob Store ────────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_storage(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "Multi-backend blob storage abstraction for documents, artifacts, and signed downloads",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # BlobRef diagram
    add_shape(s, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.5), Inches(7), Inches(2.0), fill_color=CARD_BG)
    add_shape(s, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.5), Inches(0.06), Inches(2.0), fill_color=TEAL)

    add_textbox(s, Inches(1.8), Inches(1.6), Inches(6.5), Inches(0.35),
                "BlobRef Dataclass", font_size=16, bold=True, color=WHITE)

    fields = [
        ("backend", "Storage provider (S3, Azure, local)"),
        ("uri", "Full resource URI"),
        ("bucket / key", "Bucket name + object key"),
        ("etag / sha1", "Content integrity hashes"),
        ("size / content_type", "Byte size + MIME type"),
    ]
    y = Inches(2.05)
    for field, desc in fields:
        add_textbox(s, Inches(1.9), y, Inches(2.2), Inches(0.25),
                    field, font_size=12, bold=True, color=TEAL, font_name="Consolas")
        add_textbox(s, Inches(4.2), y, Inches(4.0), Inches(0.25),
                    desc, font_size=11, color=LIGHT_GRAY)
        y += Inches(0.3)

    # Bottom: Three backend boxes
    backends = [
        ("S3", "Amazon S3 compatible", TEAL),
        ("Azure Blob", "Azure storage accounts", ACCENT_BLUE),
        ("Local FS", "File-based development", ACCENT_GREEN),
    ]
    x = Inches(1.5)
    for name, desc, color in backends:
        add_shape(s, MSO_SHAPE.RECTANGLE, x, Inches(3.85), Inches(2.2), Inches(0.7), fill_color=color)
        add_textbox(s, x + Inches(0.15), Inches(3.9), Inches(1.9), Inches(0.25),
                    name, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), Inches(4.18), Inches(1.9), Inches(0.25),
                    desc, font_size=10, color=RGBColor(0xEE, 0xEE, 0xEE), align=PP_ALIGN.CENTER)
        x += Inches(2.4)

    # Integration note
    add_textbox(s, MARGIN, Inches(4.8), Inches(8.9), Inches(0.5),
                "Used across API for uploaded documents, agent artifacts, and signed downloads (/v1/files/{download_id}). "
                "Key sanitization + SHA1 hashing ensure safe, deduplicated storage.",
                font_size=10, color=RGBColor(0x88, 0x88, 0x88))

build_content_slide(slide, "Storage & Blob Store (storage/)", body_storage)
add_notes(slide,
    "The storage subsystem at storage/blob_store.py provides a multi-backend abstraction for file storage. "
    "BlobRef is the core dataclass representing a stored object with fields for backend type, URI, bucket, key, "
    "ETag, SHA1 hash, size, and content type. "
    "Three backends are supported: Amazon S3 (and compatible), Azure Blob Storage, and local filesystem for development. "
    "The system integrates with the API's upload and download endpoints. "
    "When agents produce artifacts (charts, reports, exported files), they're stored as blobs "
    "and accessible via the /v1/files/{download_id} signed download endpoint. "
    "Key sanitization ensures safe storage paths, and SHA1 hashing enables content deduplication.")


# ─── NEW SLIDE: Graph Query Methods ─────────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_graph_methods(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "4 query methods, phased build pipeline, and 2 backend options",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # Four method cards
    methods = [
        ("LOCAL", "Entity & relationship queries. Multihop within document clusters. Requires 5 artifacts.",
         TEAL),
        ("GLOBAL", "Corpus-wide summaries and cross-document themes. Requires 3 artifacts.",
         ACCENT_ORANGE),
        ("DRIFT", "Temporal change detection and evolution tracking across rebuilds. 5 artifacts.",
         ACCENT_PURPLE),
        ("SQL", "Structured queries, exact matches, metadata filtering via StructuredSearchAdapter.",
         DARK_TEAL),
    ]
    cw = Inches(4.35)
    positions = [(MARGIN, Inches(1.5)), (Inches(5.15), Inches(1.5)),
                 (MARGIN, Inches(2.85)), (Inches(5.15), Inches(2.85))]
    for (title, body, color), (x, y) in zip(methods, positions):
        add_card(s, x, y, cw, Inches(1.15), title, body,
                 accent_color=color, title_size=13, body_size=10)

    # Bottom: Build phases
    add_textbox(s, MARGIN, Inches(4.3), Inches(9), Inches(0.35),
                "PHASED BUILD PIPELINE", font_size=12, bold=True, color=NAVY)

    phases = [
        ("Phase 1", "Entity/Relationship Extraction via LLM"),
        ("Phase 2", "Community Detection (native or text-based fallback)"),
        ("Phase 3", "Embedding & Indexing"),
    ]
    x = MARGIN
    for label, desc in phases:
        add_shape(s, MSO_SHAPE.RECTANGLE, x, Inches(4.7), Inches(2.8), Inches(0.65), fill_color=NAVY)
        add_textbox(s, x + Inches(0.15), Inches(4.75), Inches(2.5), Inches(0.22),
                    label, font_size=11, bold=True, color=TEAL)
        add_textbox(s, x + Inches(0.15), Inches(5.0), Inches(2.5), Inches(0.3),
                    desc, font_size=10, color=LIGHT_GRAY)
        x += Inches(3.05)

build_content_slide(slide, "Graph Query Methods & Build Pipeline (graph/)", body_graph_methods)
add_notes(slide,
    "The graph subsystem supports 4 distinct query methods. "
    "LOCAL queries work at the entity and relationship level, doing multihop within document clusters. "
    "GLOBAL queries produce corpus-wide summaries and cross-document theme analysis using community insights. "
    "DRIFT is notable - it does temporal change detection across graph index rebuilds, tracking how entities and relationships evolve over time. "
    "SQL queries use the StructuredSearchAdapter for exact matches and metadata filtering. "
    "Query method aliases: 'graph' maps to (local, global), 'multihop' to (local, global), 'relationship' to just (local). "
    "The build pipeline runs in 3 phases: Phase 1 extracts entities and relationships via LLM, "
    "Phase 2 detects communities (with a text-based fallback if native detection fails), "
    "Phase 3 generates embeddings and indexes. "
    "Two backends: Microsoft GraphRAG and Neo4j. Artifacts are lazy-loaded with an in-memory cache.")


# ─── NEW SLIDE: Configuration & Context Budget ──────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_config(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "300+ configuration variables in config.py with fine-grained subsystem tuning",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # Left: Context Budget
    add_textbox(s, MARGIN, Inches(1.4), Inches(4.3), Inches(0.35),
                "CONTEXT BUDGET", font_size=13, bold=True, color=TEAL)

    budget_items = [
        ("context_window_tokens", "32,768"),
        ("context_target_ratio", "0.72"),
        ("autocompact_threshold", "0.85"),
        ("tool_result_max_tokens", "2,000"),
        ("tool_results_total", "8,000"),
        ("microcompact_target", "2,400"),
        ("compact_recent_msgs", "12"),
    ]
    y = Inches(1.85)
    for key, val in budget_items:
        add_textbox(s, MARGIN, y, Inches(2.6), Inches(0.22),
                    key, font_size=10, color=RGBColor(0x44, 0x44, 0x44), font_name="Consolas")
        add_textbox(s, Inches(3.2), y, Inches(1.3), Inches(0.22),
                    val, font_size=10, bold=True, color=TEAL, font_name="Consolas")
        y += Inches(0.28)

    # Right: Worker Scheduler
    add_textbox(s, Inches(5.3), Inches(1.4), Inches(4.2), Inches(0.35),
                "WORKER SCHEDULER", font_size=13, bold=True, color=ACCENT_ORANGE)

    scheduler_items = [
        ("max_concurrency", "6"),
        ("urgent weight", "8"),
        ("interactive weight", "3"),
        ("background weight", "1"),
        ("budget_tokens/min", "24,000"),
        ("burst_tokens", "48,000"),
        ("token_cost_factor", "1.35x"),
    ]
    y = Inches(1.85)
    for key, val in scheduler_items:
        add_textbox(s, Inches(5.3), y, Inches(2.6), Inches(0.22),
                    key, font_size=10, color=RGBColor(0x44, 0x44, 0x44), font_name="Consolas")
        add_textbox(s, Inches(8.0), y, Inches(1.3), Inches(0.22),
                    val, font_size=10, bold=True, color=ACCENT_ORANGE, font_name="Consolas")
        y += Inches(0.28)

    # Bottom row: RAG + Memory defaults
    add_card(s, MARGIN, Inches(4.0), Inches(4.3), Inches(1.35),
             "RAG Defaults",
             "top_k: 15 (vector + keyword)  |  chunk: 900/150\n"
             "RRF k=60  |  retries: 2  |  min_evidence: 2\n"
             "parallel_collection_probes: 4",
             accent_color=TEAL, title_size=12, body_size=10)

    add_card(s, Inches(5.15), Inches(4.0), Inches(4.35), Inches(1.35),
             "Memory Defaults",
             "top_k: 16  |  budget: 1,600 chars\n"
             "Modes: shadow / selector / live\n"
             "5 types: profile, task, decision, constraint, open_loop",
             accent_color=ACCENT_PURPLE, title_size=12, body_size=10)

build_content_slide(slide, "Configuration & Tuning Defaults (config.py)", body_config)
add_notes(slide,
    "config.py contains over 300 configuration variables organized by subsystem. "
    "Context budget management is critical: 32,768 token window, 72% target ratio, 85% autocompact threshold. "
    "Tool results are capped at 2,000 tokens each and 8,000 total. "
    "The microcompact target of 2,400 tokens is the budget for aggressive context compression. "
    "Worker scheduler uses 3 queue classes with weighted scheduling: urgent(8), interactive(3), background(1). "
    "Per-tenant budget tracking: 24,000 tokens per minute with 48,000 burst capacity. "
    "Token cost estimation multiplies text tokens by 1.35 to account for completion overhead. "
    "RAG defaults: 15 results each for vector and keyword search, 900-char chunks with 150 overlap. "
    "Reciprocal Rank Fusion with k=60 for result merging. "
    "Memory has 5 types with weighted scoring. Shadow mode tests writes without affecting production. "
    "Episode triggers fire on 8+ messages or plan/status/decision keywords.")


# ─── NEW SLIDE: Deep RAG & Verification Pipeline ────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_deep_rag(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "Multi-pass parallel retrieval with automated verification and citation checking",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # Left: Deep RAG
    add_textbox(s, MARGIN, Inches(1.4), Inches(4.3), Inches(0.35),
                "DEEP RAG MODE", font_size=13, bold=True, color=TEAL)

    deep_items = [
        ("Parallel Lanes:", "3 concurrent retrieval lanes"),
        ("Chunk Threshold:", "24 chunks trigger full-read mode"),
        ("Reflection Rounds:", "1 sync reflection per pass"),
        ("Background Threshold:", "4+ tasks trigger background mode"),
        ("Coverage Goals:", "corpus_wide, exhaustive, specific"),
        ("Result Modes:", "answer, comparison, inventory"),
        ("Search Modes:", "auto, vector_only, keyword_only, hybrid, none"),
    ]
    y = Inches(1.85)
    for label, val in deep_items:
        add_textbox(s, MARGIN, y, Inches(1.9), Inches(0.22),
                    label, font_size=10, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, Inches(2.5), y, Inches(2.3), Inches(0.22),
                    val, font_size=10, color=RGBColor(0x55, 0x55, 0x55))
        y += Inches(0.3)

    # Right: Verification Pipeline
    add_textbox(s, Inches(5.3), Inches(1.4), Inches(4.2), Inches(0.35),
                "VERIFICATION PIPELINE", font_size=13, bold=True, color=ACCENT_RED)

    verif_items = [
        ("Missed Doc Detection", "Identifies documents that should have been retrieved but weren't"),
        ("Unsupported Hop Detection", "Flags claims requiring evidence hops the system can't make"),
        ("Citation Mismatch", "Checks citation topic/text alignment against source material"),
        ("Stale Graph Detection", "Compares graph index age vs. document ingest timestamps"),
    ]
    y = Inches(1.85)
    for title, desc in verif_items:
        add_shape(s, MSO_SHAPE.OVAL, Inches(5.3), y + Inches(0.04), Inches(0.14), Inches(0.14), fill_color=ACCENT_RED)
        add_textbox(s, Inches(5.55), y, Inches(3.9), Inches(0.22),
                    title, font_size=11, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, Inches(5.55), y + Inches(0.22), Inches(3.9), Inches(0.3),
                    desc, font_size=9, color=RGBColor(0x66, 0x66, 0x66))
        y += Inches(0.55)

    # Bottom: Relevance Grading
    add_shape(s, MSO_SHAPE.RECTANGLE, MARGIN, Inches(4.3), Inches(8.9), Inches(1.0), fill_color=NAVY)
    add_textbox(s, Inches(0.8), Inches(4.38), Inches(8.4), Inches(0.3),
                "4-LEVEL RELEVANCE GRADING", font_size=12, bold=True, color=TEAL)
    grades = "0 = Not Relevant    |    1 = Somewhat    |    2 = Relevant    |    3 = Highly Relevant"
    add_textbox(s, Inches(0.8), Inches(4.7), Inches(8.4), Inches(0.25),
                grades, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(5.0), Inches(8.4), Inches(0.22),
                "LLM grading (max 12 chunks) with heuristic fallback (term overlap >= 3 tokens = score 3)",
                font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

build_content_slide(slide, "Deep RAG & Verification Pipeline (rag/)", body_deep_rag)
add_notes(slide,
    "Deep RAG mode is a multi-pass retrieval strategy triggered for complex queries. "
    "It runs 3 parallel retrieval lanes simultaneously, with a 24-chunk threshold for switching to full-read mode. "
    "One sync reflection round per pass lets the system re-evaluate what it's found before proceeding. "
    "If 4+ retrieval tasks are needed, execution moves to background mode. "
    "Three coverage goals: corpus_wide (scan everything), exhaustive (leave no stone unturned), specific (targeted extraction). "
    "Three result modes: answer (synthesized response), comparison (side-by-side), inventory (per-document listing). "
    "The verification pipeline runs 4 checks: missed document detection identifies overlooked sources, "
    "unsupported hop detection flags claims needing evidence the system can't reach, "
    "citation mismatch checking verifies topic and text alignment, "
    "and stale graph detection compares graph age against document ingest timestamps. "
    "Relevance grading uses a 4-level scale (0-3) with LLM grading for up to 12 chunks "
    "and heuristic fallback for the rest.")


# ─── NEW SLIDE: Persistence & Observability ──────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_persistence(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "PostgreSQL persistence layer, dedicated observability directory, and benchmark suite",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66))

    # Three columns
    col_w = Inches(2.8)
    gap = Inches(0.15)

    # Column 1: Persistence
    x1 = MARGIN
    add_shape(s, MSO_SHAPE.RECTANGLE, x1, Inches(1.5), col_w, Inches(3.3), fill_color=CARD_BG)
    add_shape(s, MSO_SHAPE.RECTANGLE, x1, Inches(1.5), Inches(0.06), Inches(3.3), fill_color=TEAL)
    add_textbox(s, x1 + Inches(0.2), Inches(1.6), col_w - Inches(0.3), Inches(0.3),
                "Persistence (persistence/)", font_size=14, bold=True, color=WHITE)
    items_p = [
        "PostgreSQL-backed data layer",
        "pgvector for embeddings",
        "Session transcript storage",
        "Job records & results",
        "Tenant/user data isolation",
        "Migration support",
    ]
    y = Inches(2.0)
    for item in items_p:
        add_textbox(s, x1 + Inches(0.25), y, col_w - Inches(0.4), Inches(0.25),
                    "- " + item, font_size=10, color=LIGHT_GRAY)
        y += Inches(0.3)

    # Column 2: Observability
    x2 = x1 + col_w + gap
    add_shape(s, MSO_SHAPE.RECTANGLE, x2, Inches(1.5), col_w, Inches(3.3), fill_color=CARD_BG)
    add_shape(s, MSO_SHAPE.RECTANGLE, x2, Inches(1.5), Inches(0.06), Inches(3.3), fill_color=ACCENT_ORANGE)
    add_textbox(s, x2 + Inches(0.2), Inches(1.6), col_w - Inches(0.3), Inches(0.3),
                "Observability (observability/)", font_size=14, bold=True, color=WHITE)
    items_o = [
        "LangFuse integration for traces",
        "Event emission per tool call",
        "Router decision logging",
        "Turn-level status tracking",
        "SSE progress streaming",
        "events.jsonl audit trail",
    ]
    y = Inches(2.0)
    for item in items_o:
        add_textbox(s, x2 + Inches(0.25), y, col_w - Inches(0.4), Inches(0.25),
                    "- " + item, font_size=10, color=LIGHT_GRAY)
        y += Inches(0.3)

    # Column 3: Benchmark
    x3 = x2 + col_w + gap
    add_shape(s, MSO_SHAPE.RECTANGLE, x3, Inches(1.5), col_w, Inches(3.3), fill_color=CARD_BG)
    add_shape(s, MSO_SHAPE.RECTANGLE, x3, Inches(1.5), Inches(0.06), Inches(3.3), fill_color=ACCENT_PURPLE)
    add_textbox(s, x3 + Inches(0.2), Inches(1.6), col_w - Inches(0.3), Inches(0.3),
                "Benchmark (benchmark/)", font_size=14, bold=True, color=WHITE)
    items_b = [
        "Performance benchmarking suite",
        "Defense RAG test corpus",
        "Adversarial prompt testing",
        "Regression detection",
        "Latency profiling per tool",
        "Token cost analysis",
    ]
    y = Inches(2.0)
    for item in items_b:
        add_textbox(s, x3 + Inches(0.25), y, col_w - Inches(0.4), Inches(0.25),
                    "- " + item, font_size=10, color=LIGHT_GRAY)
        y += Inches(0.3)

    # Bottom
    add_textbox(s, MARGIN, Inches(5.0), Inches(8.9), Inches(0.4),
                "Sandbox: 180s timeout, isolated Docker containers. Control Panel Frontend: React/TypeScript (Vite) at /control_panel/.",
                font_size=10, color=RGBColor(0x88, 0x88, 0x88))

build_content_slide(slide, "Persistence, Observability & Benchmarks", body_persistence)
add_notes(slide,
    "Three supporting subsystems that round out the production infrastructure. "
    "The persistence layer uses PostgreSQL with pgvector for embedding storage. "
    "It handles session transcripts, job records, tenant and user data isolation, and migration support. "
    "The observability directory integrates with LangFuse for distributed traces and spans. "
    "Events are emitted per tool call, router decision, and agent turn. "
    "Turn-level status tracking feeds the SSE progress stream so clients see real-time updates. "
    "events.jsonl provides a local audit trail of all system activity. "
    "The benchmark suite includes a defense RAG test corpus for adversarial testing, "
    "latency profiling per tool, and token cost analysis. "
    "Additional infrastructure: sandbox timeout is 180 seconds, "
    "and the control panel frontend is a separate React/TypeScript app built with Vite at /control_panel/.")


# ══════════════════════════════════════════════════════════════
# STEP 3: REORDER — Insert new slides after slide 11
# (After "Document Processing & New Modules")
# ══════════════════════════════════════════════════════════════

print(f"\nAdded {len(new_slides)} new slides")
print("Reordering slides...")

# Get the presentation XML element
pres_elem = prs.part._element
sldIdLst = pres_elem.find(qn('p:sldIdLst'))
sldId_elements = list(sldIdLst)

# The new slides are at the end. We want to insert them after the existing
# "Document Processing & New Modules" slide (originally slide 11, now index 10 in 0-based).
# But since we've already added intro slides in the previous build, slide 11 is at position 10.

# Find where to insert: after the current slide 11 (Document Processing & New Modules)
# In the current 140-slide deck, that's index 10 in the sldId list.
# But let's be safe and insert right before the enterprise roadmap section.
# Looking at content: slide 134 is "Enterprise-Grade Roadmap", so we'll insert before that.

# Actually, let's place the new audit slides after slide 11 (Document Processing & New Modules)
# The existing deck has 140 slides. New slides are at indices 140-150.
# We want to move them to after index 10 (which is slide 11 in 1-based numbering).

total_existing = 140
num_new = len(new_slides)

# Remove the new sldId elements from the end
new_sldIds = sldId_elements[-num_new:]
for elem in new_sldIds:
    sldIdLst.remove(elem)

# Re-read after removal
sldId_elements = list(sldIdLst)

# Insert after position 10 (slide 11 = Document Processing & New Modules)
insert_after_idx = 10
for i, new_elem in enumerate(new_sldIds):
    sldIdLst.insert(insert_after_idx + 1 + i, new_elem)

print(f"Inserted {num_new} slides after slide 11")


# ══════════════════════════════════════════════════════════════
# STEP 4: UPDATE SPEAKER NOTES ON KEY EXISTING SLIDES
# ══════════════════════════════════════════════════════════════

print("\nUpdating speaker notes on existing slides...")

# We'll update notes on slides that need corrections.
# After reordering, the original slide positions shift, but notes update by content.
# Let's target specific slides by iterating and matching content.

slides_list = list(prs.slides)

for i, slide in enumerate(slides_list):
    # Find slides by their title text
    slide_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                slide_text += p.text + " "

    # Update the API endpoint slide notes
    if "API Surface" in slide_text and "118 Endpoints" in slide_text:
        set_notes(slide,
            "The API surface has grown to 118 total endpoints: 48 public and 70 admin. "
            "Public endpoints cover chat, documents, graphs, skills, MCP connections, jobs/tasks, team mailbox, "
            "files/downloads, sessions, capabilities, and system health. "
            "The admin API handles agents, collections, prompts, graphs, access control, "
            "config, uploads, MCP, architecture snapshots, and operations. "
            "All public endpoints are in api/main.py. Admin endpoints are in control_panel/routes.py.")

    # Update the "All 48 Public Endpoints" slide notes
    if "All 48 Public Endpoints" in slide_text:
        set_notes(slide,
            "48 public endpoints across 12 categories. Notable additions since v2: "
            "/v1/connector/chat for third-party integrations, /v1/agents returns agent configs, "
            "five graph endpoints (list, get, index, import, query), "
            "eight MCP connection endpoints, five jobs/tasks endpoints, "
            "four team mailbox endpoints, health probes at /health/live and /health/ready, "
            "and capabilities/users endpoints. Admin endpoints (70 of them) are in control_panel/routes.py behind X-Admin-Token auth.")

    # Update agent overview slide with corrected tool counts
    if "All 11 Agent Roles" in slide_text:
        set_notes(slide,
            "System has 11 agents. Key tool count updates from the latest audit: "
            "general has 33 tools (including MCP wildcard for dynamic tool discovery), "
            "coordinator has 11 tools (spawn_worker, message_worker, team mailbox tools, etc.), "
            "data_analyst has 20 tools (analyst + team mailbox + invoke_agent), "
            "graph_manager has 12 tools, utility has 13 tools, verifier has 13 tools (RAG + graph). "
            "basic, rag_worker, planner, and finalizer have 0 tools. memory_maintainer has 3. "
            "Coordinator now has 9 allowed workers including graph_manager.")

    # Update the "at a Glance" slide
    if "at a Glance" in slide_text and "11" in slide_text and "54" in slide_text:
        set_notes(slide,
            "Updated numbers: 11 specialized agents, 54 registered tools, 7 execution modes, 8 tool groups. "
            "Agents are markdown definitions, hot-reloadable via the 70-endpoint admin API. "
            "Three-layer runtime: RuntimeService to RuntimeKernel to QueryLoop keeps orchestration clean. "
            "GraphRAG with 8 tools and MCP for dynamic tool discovery are major additions. "
            "Worker orchestration provides 14 tools for spawn, messaging, team mailbox, and multi-agent coordination. "
            "New since last review: authorization subsystem, documents deep processing, skills telemetry and dependency graphs, "
            "provider circuit breaker, task plan artifact handoff, MCP security model, storage blob store, "
            "deep RAG mode, and 4 graph query methods including drift for temporal analysis.")

print("Speaker notes updated.")


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════

prs.save(OUT)
print(f"\nSaved to: {OUT}")
print(f"Total slides: {len(prs.slides)}")
