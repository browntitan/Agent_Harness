"""
Build updated PowerPoint deck:
1. Add 7 new intro slides explaining agent harnesses (after slide 1)
2. Update existing architecture slides to match actual codebase
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
from lxml import etree
import os

# ─── CONSTANTS ───────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
TEAL      = RGBColor(0x00, 0x97, 0xA7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEAL = RGBColor(0x00, 0x6D, 0x77)
GOLD      = RGBColor(0xFF, 0xB7, 0x4D)
MED_GRAY  = RGBColor(0x99, 0x99, 0x99)
BG_LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)
CARD_BG   = RGBColor(0x23, 0x3B, 0x5E)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
ACCENT_PURPLE = RGBColor(0x7C, 0x4D, 0xFF)
ACCENT_RED = RGBColor(0xEF, 0x53, 0x50)

SW = Inches(10)
SH = Inches(5.625)
MARGIN = Inches(0.55)

def set_bg(slide, color):
    """Set slide background to solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                font_name="Arial", valign=MSO_ANCHOR.TOP, line_spacing=None):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    tf.auto_size = None
    # Set vertical alignment
    txBox.text_frame._txBody.bodyPr.set("anchor", {
        MSO_ANCHOR.TOP: "t",
        MSO_ANCHOR.MIDDLE: "ctr",
        MSO_ANCHOR.BOTTOM: "b"
    }.get(valign, "t"))
    if line_spacing:
        from pptx.oxml.ns import qn
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txBox

def add_rich_textbox(slide, left, top, width, height, runs, line_spacing=None):
    """Add textbox with multiple formatted runs.
    runs: list of dicts with keys: text, font_size, bold, color, font_name, breakLine, align, bullet
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    current_p = tf.paragraphs[0]
    first_run = True

    for run_def in runs:
        if run_def.get('breakLine') and not first_run:
            current_p = tf.add_paragraph()

        if run_def.get('align'):
            current_p.alignment = run_def['align']

        r = current_p.add_run()
        r.text = run_def.get('text', '')
        r.font.size = Pt(run_def.get('font_size', 14))
        r.font.bold = run_def.get('bold', False)
        r.font.color.rgb = run_def.get('color', WHITE)
        r.font.name = run_def.get('font_name', 'Arial')

        if line_spacing:
            from pptx.oxml.ns import qn
            pPr = current_p._pPr if current_p._pPr is not None else current_p._p.get_or_add_pPr()
            if pPr.find(qn('a:lnSpc')) is None:
                lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
                spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
                spcPts.set('val', str(int(line_spacing * 100)))

        first_run = False

    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a shape with optional fill and line."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body, accent_color=TEAL,
             title_size=16, body_size=12, bg_color=CARD_BG):
    """Add a card with accent bar, title, and body text."""
    # Card background
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill_color=bg_color)
    # Accent bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height, fill_color=accent_color)
    # Title
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.3), Inches(0.4),
                title, font_size=title_size, bold=True, color=WHITE)
    # Body
    add_textbox(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6),
                body, font_size=body_size, color=LIGHT_GRAY, line_spacing=18)

def build_section_divider(slide, title, subtitle):
    """Build a section divider slide matching existing style."""
    set_bg(slide, NAVY)
    # Teal accent bar at top
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, Inches(0.06), fill_color=TEAL)
    # Title centered
    add_textbox(slide, MARGIN, Inches(1.8), Inches(9), Inches(1.0),
                title, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Subtitle
    add_textbox(slide, MARGIN, Inches(2.8), Inches(8), Inches(0.8),
                subtitle, font_size=16, color=TEAL, align=PP_ALIGN.LEFT)

def build_content_slide(slide, title, body_callback):
    """Build a content slide with title bar and body area."""
    set_bg(slide, RGBColor(0xF5, 0xF7, 0xFA))
    # Title bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, Inches(0.75), fill_color=NAVY)
    add_textbox(slide, MARGIN, Inches(0.12), Inches(9), Inches(0.5),
                title, font_size=22, bold=True, color=WHITE)
    # Teal accent line under title
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.75), SW, Inches(0.04), fill_color=TEAL)
    # Call body builder
    body_callback(slide)


# ══════════════════════════════════════════════════════════════
# MAIN BUILD
# ══════════════════════════════════════════════════════════════

SRC = "agentic_chatbot_v3/claude_powerpoint/agentic_chatbot_v2_deep_dive  -  Repaired.pptx"
OUT = "agentic_chatbot_v3/claude_powerpoint/agentic_chatbot_v2_deep_dive_updated.pptx"

prs = Presentation(SRC)
layout = prs.slide_layouts[0]

# We'll insert new slides after slide index 0 (title slide)
# python-pptx doesn't have insert_slide, so we add at end and reorder XML

new_slides = []

# ─── INTRO SLIDE 1: Section Divider — "What Is an Agent Harness?" ──────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)
build_section_divider(slide,
    "What Is an Agent Harness?",
    "Understanding the orchestration layer that turns LLMs into autonomous agents")

# ─── INTRO SLIDE 2: The Problem — Why LLMs Alone Aren't Enough ─────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_problem(s):
    # Left column — The Problem
    add_textbox(s, MARGIN, Inches(1.0), Inches(4.2), Inches(0.4),
                "THE PROBLEM", font_size=14, bold=True, color=TEAL)

    problems = [
        ("Single-Turn Only", "Raw LLMs answer one question at a time — no actions, retrieval, or state."),
        ("No Tool Access", "LLMs cannot search databases, execute code, or call APIs on their own."),
        ("No Memory", "Each request is stateless — the model forgets everything between calls."),
        ("No Planning", "Complex multi-step tasks with parallel work or verification are impossible."),
    ]
    y = Inches(1.45)
    for title, desc in problems:
        add_card(s, MARGIN, y, Inches(4.2), Inches(0.85), title, desc,
                accent_color=ACCENT_RED, bg_color=RGBColor(0x2C, 0x1F, 0x1F),
                title_size=13, body_size=10)
        y += Inches(0.92)

    # Right column — The Solution
    add_textbox(s, Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.4),
                "THE SOLUTION: AGENT HARNESS", font_size=14, bold=True, color=ACCENT_GREEN)

    add_textbox(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(3.6),
                "", font_size=12, color=LIGHT_GRAY)

    solutions = [
        ("Multi-Turn Loops", "The harness wraps the LLM in a reasoning loop — observe, think, act, repeat."),
        ("Tool Integration", "A registry of callable tools (search, code exec, APIs) invoked by name."),
        ("Session & Memory", "Persistent state across turns with long-term memory extraction and injection."),
        ("Task Orchestration", "Planners decompose complex work into subtasks dispatched in parallel."),
    ]
    y = Inches(1.45)
    for title, desc in solutions:
        add_card(s, Inches(5.2), y, Inches(4.3), Inches(0.85), title, desc,
                accent_color=ACCENT_GREEN, bg_color=RGBColor(0x1F, 0x2C, 0x1F),
                title_size=13, body_size=10)
        y += Inches(0.92)

build_content_slide(slide, "Why LLMs Alone Aren't Enough", body_problem)

# ─── INTRO SLIDE 3: Anatomy of an Agent Harness ──────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_anatomy(s):
    # Central flow diagram using shapes
    bw = Inches(1.55)
    bh = Inches(0.7)
    arrow_w = Inches(0.3)

    # Top label
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "THE CORE LOOP: Observe \u2192 Think \u2192 Act \u2192 Repeat",
                font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Flow boxes
    components = [
        ("User\nInput", TEAL),
        ("Router /\nClassifier", DARK_TEAL),
        ("Agent\n(ReAct Loop)", RGBColor(0x00, 0x6D, 0x77)),
        ("Tool\nExecution", ACCENT_ORANGE),
        ("Response\nStream", ACCENT_GREEN),
    ]

    y_top = Inches(1.5)
    positions = []
    for i, (label, color) in enumerate(components):
        x = Inches(0.35) + i * (bw + arrow_w)
        shape = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y_top, bw, bh, fill_color=color)
        shape.text_frame.word_wrap = True
        p = shape.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        positions.append((x, bw))

        # Arrow between boxes
        if i < len(components) - 1:
            ax = x + bw + Inches(0.02)
            add_textbox(s, ax, y_top + Inches(0.15), Inches(0.25), Inches(0.4),
                       "\u25B6", font_size=16, color=NAVY, align=PP_ALIGN.CENTER)

    # Bottom section: 5 key layers
    add_textbox(s, MARGIN, Inches(2.5), Inches(9), Inches(0.35),
                "KEY LAYERS OF THE HARNESS", font_size=14, bold=True, color=NAVY)

    layers = [
        ("1. Routing Layer", "Classifies intent, selects the right agent", TEAL),
        ("2. Agent Registry", "Declarative agent definitions (mode, tools, prompt)", DARK_TEAL),
        ("3. Execution Engine", "ReAct loop, RAG pipeline, coordinator pattern", RGBColor(0x00, 0x6D, 0x77)),
        ("4. Tool System", "Registry of callable functions with policy enforcement", ACCENT_ORANGE),
        ("5. Memory & State", "Session persistence, long-term memory, context budgets", ACCENT_GREEN),
    ]

    y = Inches(2.9)
    for label, desc, color in layers:
        # Accent dot
        add_shape(s, MSO_SHAPE.OVAL, MARGIN, y + Inches(0.08), Inches(0.18), Inches(0.18), fill_color=color)
        add_textbox(s, MARGIN + Inches(0.25), y, Inches(2.2), Inches(0.35),
                   label, font_size=12, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, Inches(3.0), y, Inches(6.5), Inches(0.35),
                   desc, font_size=11, color=RGBColor(0x66, 0x66, 0x66))
        y += Inches(0.42)

build_content_slide(slide, "Anatomy of an Agent Harness", body_anatomy)

# ─── INTRO SLIDE 4: The ReAct Pattern ─────────────────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_react(s):
    # Left side — ReAct loop diagram
    add_textbox(s, MARGIN, Inches(0.9), Inches(4.5), Inches(0.35),
                "THE REACT LOOP (Reasoning + Acting)", font_size=14, bold=True, color=NAVY)

    # Circular flow representation
    steps = [
        ("1. OBSERVE", "Receive user query +\nconversation context", TEAL),
        ("2. THINK", "LLM reasons about what\nto do next (chain-of-thought)", RGBColor(0x5C, 0x6B, 0xC0)),
        ("3. ACT", "Call a tool (search, code exec,\nAPI call) or respond", ACCENT_ORANGE),
        ("4. EVALUATE", "Check result \u2014 is the task\ncomplete? Loop or finish.", ACCENT_GREEN),
    ]

    y = Inches(1.4)
    for title, desc, color in steps:
        add_shape(s, MSO_SHAPE.RECTANGLE, MARGIN, y, Inches(0.08), Inches(0.75), fill_color=color)
        add_textbox(s, MARGIN + Inches(0.2), y, Inches(4.0), Inches(0.3),
                   title, font_size=13, bold=True, color=color)
        add_textbox(s, MARGIN + Inches(0.2), y + Inches(0.3), Inches(4.0), Inches(0.45),
                   desc, font_size=10, color=RGBColor(0x55, 0x55, 0x55))
        y += Inches(0.85)

    # Right side — example trace
    add_textbox(s, Inches(5.3), Inches(0.9), Inches(4.2), Inches(0.35),
                "EXAMPLE: \"What does the policy say about remote work?\"",
                font_size=12, bold=True, color=NAVY)

    # Example trace card
    add_shape(s, MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(1.35), Inches(4.2), Inches(3.8),
              fill_color=RGBColor(0x1B, 0x2A, 0x4A))

    trace_lines = [
        ("\u25B6 Router:", " RAG intent detected (0.92 confidence)", TEAL, WHITE),
        ("\u25B6 Agent:", " rag_worker selected", TEAL, WHITE),
        ("", "", None, None),
        ("Step 1 \u2014 Observe:", " Parse query for search terms", ACCENT_ORANGE, LIGHT_GRAY),
        ("Step 2 \u2014 Act:", " Hybrid search (vector + keyword)", ACCENT_ORANGE, LIGHT_GRAY),
        ("", "  \u2192 Found 12 chunks from HR_Policy.pdf", None, MED_GRAY),
        ("Step 3 \u2014 Think:", " Grade relevance of chunks", ACCENT_ORANGE, LIGHT_GRAY),
        ("", "  \u2192 5 chunks scored 2+ (relevant)", None, MED_GRAY),
        ("Step 4 \u2014 Act:", " Synthesize grounded answer", ACCENT_ORANGE, LIGHT_GRAY),
        ("", "  \u2192 Answer with 3 citations", None, MED_GRAY),
        ("", "", None, None),
        ("\u2713 Response:", " Streamed via SSE with citations", ACCENT_GREEN, WHITE),
    ]

    y = Inches(1.5)
    for label, value, lcolor, vcolor in trace_lines:
        if label == "" and value == "":
            y += Inches(0.15)
            continue
        runs = []
        if label:
            runs.append({'text': label, 'font_size': 10, 'bold': True, 'color': lcolor or WHITE})
        if value:
            runs.append({'text': value, 'font_size': 10, 'bold': False, 'color': vcolor or LIGHT_GRAY})
        add_rich_textbox(s, Inches(5.5), y, Inches(3.8), Inches(0.25), runs)
        y += Inches(0.28)

build_content_slide(slide, "The ReAct Pattern — How Agents Reason", body_react)

# ─── INTRO SLIDE 5: Agent Modes & Specialization ──────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_modes(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.35),
                "Different tasks need different execution strategies. The harness routes to specialized agent modes.",
                font_size=13, color=RGBColor(0x55, 0x55, 0x55))

    modes = [
        ("BASIC", "Single LLM call, no tools", "Fast Q&A, greetings, simple tasks", TEAL, "//"),
        ("REACT", "ReAct loop with tool calling", "General tasks, data analysis, file ops", RGBColor(0x5C, 0x6B, 0xC0), "<>"),
        ("RAG", "Retrieve > Grade > Synthesize", "Document search, policy questions", ACCENT_ORANGE, "{}"),
        ("COORD", "Planner > Workers > Finalizer", "Complex multi-step research", ACCENT_PURPLE, ">>"),
    ]

    col_w = Inches(2.15)
    x_start = Inches(0.3)
    y_top = Inches(1.5)

    for i, (name, method, example, color, icon) in enumerate(modes):
        x = x_start + i * (col_w + Inches(0.15))
        # Card
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y_top, col_w, Inches(3.4), fill_color=NAVY)
        # Accent top
        add_shape(s, MSO_SHAPE.RECTANGLE, x, y_top, col_w, Inches(0.06), fill_color=color)
        # Icon
        add_textbox(s, x, y_top + Inches(0.2), col_w, Inches(0.5),
                   icon, font_size=28, color=color, align=PP_ALIGN.CENTER)
        # Mode name
        add_textbox(s, x + Inches(0.15), y_top + Inches(0.75), col_w - Inches(0.3), Inches(0.35),
                   name, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Method
        add_textbox(s, x + Inches(0.15), y_top + Inches(1.2), col_w - Inches(0.3), Inches(0.6),
                   method, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        # Divider
        add_shape(s, MSO_SHAPE.RECTANGLE, x + Inches(0.3), y_top + Inches(1.9),
                 col_w - Inches(0.6), Inches(0.02), fill_color=color)
        # Example
        add_textbox(s, x + Inches(0.15), y_top + Inches(2.05), col_w - Inches(0.3), Inches(1.0),
                   example, font_size=10, color=MED_GRAY, align=PP_ALIGN.CENTER)

build_content_slide(slide, "Agent Modes — One Harness, Many Strategies", body_modes)

# ─── INTRO SLIDE 6: Tool System & RAG Overview ───────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_tools(s):
    # Left: Tool system
    add_textbox(s, MARGIN, Inches(0.9), Inches(4.3), Inches(0.35),
                "TOOL SYSTEM", font_size=14, bold=True, color=NAVY)
    add_textbox(s, MARGIN, Inches(1.25), Inches(4.3), Inches(0.5),
                "54 tools in 8 groups, registered with the harness. Each agent gets a curated subset via its allowed_tools whitelist.",
                font_size=11, color=RGBColor(0x66, 0x66, 0x66))

    tool_groups = [
        ("Orchestration (14)", "spawn/message workers, team mailbox, invoke agent, job control", ACCENT_PURPLE),
        ("RAG Gateway (12)", "search, read, compare docs, extract, consolidate, evidence binder", TEAL),
        ("Analyst (11)", "load_dataset, execute_code (Docker), scratchpad, workspace ops", ACCENT_ORANGE),
        ("Graph Gateway (8)", "index/search/import graphs, explain source plan", RGBColor(0x5C, 0x6B, 0xC0)),
        ("Memory (3)", "memory_save, memory_load, memory_list", ACCENT_GREEN),
        ("Utility (3)", "calculator, list_indexed_docs, search_skills", MED_GRAY),
        ("Discovery (2)", "discover_tools, call_deferred_tool", DARK_TEAL),
        ("Skills (1)", "execute_skill (dynamic skill execution)", RGBColor(0x88, 0x88, 0x88)),
    ]

    y = Inches(1.85)
    for name, desc, color in tool_groups:
        add_shape(s, MSO_SHAPE.OVAL, MARGIN, y + Inches(0.03), Inches(0.14), Inches(0.14), fill_color=color)
        add_textbox(s, MARGIN + Inches(0.22), y, Inches(1.7), Inches(0.22),
                   name, font_size=10, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, MARGIN + Inches(0.22), y + Inches(0.19), Inches(4.0), Inches(0.3),
                   desc, font_size=8, color=RGBColor(0x77, 0x77, 0x77))
        y += Inches(0.47)

    # Right: RAG pipeline
    add_textbox(s, Inches(5.2), Inches(0.9), Inches(4.3), Inches(0.35),
                "RAG PIPELINE", font_size=14, bold=True, color=NAVY)
    add_textbox(s, Inches(5.2), Inches(1.25), Inches(4.3), Inches(0.5),
                "Retrieval-Augmented Generation grounds answers in your actual documents with citations.",
                font_size=11, color=RGBColor(0x66, 0x66, 0x66))

    rag_steps = [
        ("1", "Query Analysis", "Parse intent, extract entities", TEAL),
        ("2", "Hybrid Retrieval", "Vector + keyword search in parallel", RGBColor(0x5C, 0x6B, 0xC0)),
        ("3", "Judge Grading", "LLM scores each chunk 0-3 for relevance", ACCENT_ORANGE),
        ("4", "Evidence Assembly", "Deduplicate, rank, budget context", ACCENT_PURPLE),
        ("5", "Synthesis", "Generate grounded answer with citations", ACCENT_GREEN),
    ]

    y = Inches(1.85)
    for num, title, desc, color in rag_steps:
        # Number circle
        circ = add_shape(s, MSO_SHAPE.OVAL, Inches(5.2), y, Inches(0.3), Inches(0.3), fill_color=color)
        circ.text_frame.paragraphs[0].text = num
        circ.text_frame.paragraphs[0].font.size = Pt(11)
        circ.text_frame.paragraphs[0].font.bold = True
        circ.text_frame.paragraphs[0].font.color.rgb = WHITE
        circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(s, Inches(5.6), y, Inches(1.8), Inches(0.3),
                   title, font_size=11, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, Inches(7.3), y, Inches(2.3), Inches(0.3),
                   desc, font_size=10, color=RGBColor(0x66, 0x66, 0x66))

        # Connector line
        if num != "5":
            add_shape(s, MSO_SHAPE.RECTANGLE, Inches(5.34), y + Inches(0.32),
                     Inches(0.02), Inches(0.25), fill_color=RGBColor(0xDD, 0xDD, 0xDD))
        y += Inches(0.55)

build_content_slide(slide, "Tools & RAG — The Agent's Capabilities", body_tools)

# ─── INTRO SLIDE 7: Our System at a Glance ──────────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_our_system(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.4),
                "This deck deep-dives into our production agent harness. Here's what makes it distinctive:",
                font_size=13, color=RGBColor(0x55, 0x55, 0x55))

    stats = [
        ("11", "Specialized\nAgents", TEAL),
        ("54", "Registered\nTools", ACCENT_ORANGE),
        ("7", "Execution\nModes", ACCENT_PURPLE),
        ("8", "Tool\nGroups", ACCENT_GREEN),
    ]

    # Stats row
    stat_w = Inches(1.8)
    x_start = Inches(0.7)
    for i, (number, label, color) in enumerate(stats):
        x = x_start + i * (stat_w + Inches(0.45))
        add_textbox(s, x, Inches(1.5), stat_w, Inches(0.7),
                   number, font_size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(2.2), stat_w, Inches(0.5),
                   label, font_size=12, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)

    # Key differentiators
    add_shape(s, MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), SW, Inches(0.04), fill_color=RGBColor(0xDD, 0xDD, 0xDD))

    differentiators = [
        ("Agents as Data", "Markdown definitions, not hard-coded classes \u2014 hot-reloadable via 70-endpoint admin API"),
        ("Three-Layer Runtime", "RuntimeService \u2192 RuntimeKernel \u2192 QueryLoop separation keeps orchestration clean"),
        ("GraphRAG + MCP", "Knowledge graph integration (8 tools) and Model Context Protocol for dynamic tool discovery"),
        ("Worker Orchestration", "14 orchestration tools: spawn workers, team mailbox, multi-agent coordination"),
    ]

    y = Inches(3.2)
    for i, (title, desc) in enumerate(differentiators):
        col = Inches(0.55) if i % 2 == 0 else Inches(5.2)
        add_shape(s, MSO_SHAPE.RECTANGLE, col, y, Inches(4.3), Inches(0.02), fill_color=TEAL)
        add_textbox(s, col, y + Inches(0.08), Inches(4.3), Inches(0.3),
                   title, font_size=13, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, col, y + Inches(0.35), Inches(4.3), Inches(0.5),
                   desc, font_size=10, color=RGBColor(0x66, 0x66, 0x66))
        if i == 1:
            y += Inches(0.85)

build_content_slide(slide, "Our Agent Harness at a Glance", body_our_system)


# ─── INTRO SLIDE 8: New Subsystems — GraphRAG, MCP, Worker Orchestration ──
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_new_subsystems(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.35),
                "Three major subsystems added since v2 that expand the harness well beyond simple tool calling.",
                font_size=12, color=RGBColor(0x55, 0x55, 0x55))

    # Column 1: GraphRAG
    c1x = MARGIN
    cw = Inches(2.85)
    add_shape(s, MSO_SHAPE.RECTANGLE, c1x, Inches(1.45), cw, Inches(3.7), fill_color=NAVY)
    add_shape(s, MSO_SHAPE.RECTANGLE, c1x, Inches(1.45), cw, Inches(0.06), fill_color=TEAL)
    add_textbox(s, c1x + Inches(0.15), Inches(1.6), cw - Inches(0.3), Inches(0.35),
               "GRAPHRAG", font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, c1x + Inches(0.15), Inches(2.0), cw - Inches(0.3), Inches(0.3),
               "Knowledge Graph Integration", font_size=11, color=TEAL, align=PP_ALIGN.CENTER)
    add_shape(s, MSO_SHAPE.RECTANGLE, c1x + Inches(0.4), Inches(2.4), cw - Inches(0.8), Inches(0.02), fill_color=RGBColor(0x33, 0x55, 0x77))
    graph_items = "graph/ module with 10 files\ngraph_manager agent\n8 dedicated tools\nPhased corpus indexing\nCommunity report recovery\nStructured search + planning\nPrompt tuning per graph"
    add_textbox(s, c1x + Inches(0.15), Inches(2.55), cw - Inches(0.3), Inches(2.5),
               graph_items, font_size=9, color=LIGHT_GRAY, line_spacing=16)

    # Column 2: MCP
    c2x = Inches(3.55)
    add_shape(s, MSO_SHAPE.RECTANGLE, c2x, Inches(1.45), cw, Inches(3.7), fill_color=NAVY)
    add_shape(s, MSO_SHAPE.RECTANGLE, c2x, Inches(1.45), cw, Inches(0.06), fill_color=ACCENT_ORANGE)
    add_textbox(s, c2x + Inches(0.15), Inches(1.6), cw - Inches(0.3), Inches(0.35),
               "MCP", font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, c2x + Inches(0.15), Inches(2.0), cw - Inches(0.3), Inches(0.3),
               "Model Context Protocol", font_size=11, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
    add_shape(s, MSO_SHAPE.RECTANGLE, c2x + Inches(0.4), Inches(2.4), cw - Inches(0.8), Inches(0.02), fill_color=RGBColor(0x55, 0x44, 0x33))
    mcp_items = "mcp/ service + client modules\nDynamic tool registration\n8 API endpoints for connections\nSecurity-scoped tool access\nConnection test + refresh\nPer-tool enable/disable\nExternal server integration"
    add_textbox(s, c2x + Inches(0.15), Inches(2.55), cw - Inches(0.3), Inches(2.5),
               mcp_items, font_size=9, color=LIGHT_GRAY, line_spacing=16)

    # Column 3: Worker Orchestration
    c3x = Inches(6.6)
    add_shape(s, MSO_SHAPE.RECTANGLE, c3x, Inches(1.45), cw, Inches(3.7), fill_color=NAVY)
    add_shape(s, MSO_SHAPE.RECTANGLE, c3x, Inches(1.45), cw, Inches(0.06), fill_color=ACCENT_PURPLE)
    add_textbox(s, c3x + Inches(0.15), Inches(1.6), cw - Inches(0.3), Inches(0.35),
               "WORKER SYSTEM", font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, c3x + Inches(0.15), Inches(2.0), cw - Inches(0.3), Inches(0.3),
               "Multi-Agent Coordination", font_size=11, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)
    add_shape(s, MSO_SHAPE.RECTANGLE, c3x + Inches(0.4), Inches(2.4), cw - Inches(0.8), Inches(0.02), fill_color=RGBColor(0x44, 0x33, 0x55))
    worker_items = "14 orchestration tools\nspawn_worker / message_worker\nTeam mailbox channels\nParent question/approval flow\nWorker request/response cycle\nJob listing and control\nCross-agent artifact handoff"
    add_textbox(s, c3x + Inches(0.15), Inches(2.55), cw - Inches(0.3), Inches(2.5),
               worker_items, font_size=9, color=LIGHT_GRAY, line_spacing=16)

build_content_slide(slide, "New Subsystems — GraphRAG, MCP, Worker Orchestration", body_new_subsystems)


# ─── INTRO SLIDE 9: API Surface ──────────────────────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_api(s):
    add_textbox(s, MARGIN, Inches(0.9), Inches(9), Inches(0.35),
                "47 public endpoints + 70 admin endpoints = 117 total API surface",
                font_size=13, bold=True, color=RGBColor(0x44, 0x44, 0x44))

    # Left: Public API
    add_textbox(s, MARGIN, Inches(1.4), Inches(4.3), Inches(0.35),
                "PUBLIC API (47 endpoints)", font_size=14, bold=True, color=TEAL)

    pub_groups = [
        ("Chat", "/v1/chat/completions, /v1/connector/chat", "2"),
        ("Documents", "upload, ingest, source download", "4"),
        ("Graphs", "list, get, index, import, query", "5"),
        ("Skills", "CRUD, activate, deactivate, preview, rollback", "9"),
        ("MCP", "connections CRUD, tools, test, refresh", "8"),
        ("Jobs/Tasks", "get, mailbox, respond, list, stop", "6"),
        ("Team Mailbox", "channels, messages, respond", "5"),
        ("System", "health, models, agents, diagnostics, capabilities", "8"),
    ]

    y = Inches(1.8)
    for name, desc, count in pub_groups:
        add_shape(s, MSO_SHAPE.RECTANGLE, MARGIN, y, Inches(0.06), Inches(0.4), fill_color=TEAL)
        add_textbox(s, MARGIN + Inches(0.15), y, Inches(1.3), Inches(0.22),
                   name, font_size=10, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, MARGIN + Inches(0.15), y + Inches(0.2), Inches(3.5), Inches(0.2),
                   desc, font_size=8, color=RGBColor(0x77, 0x77, 0x77))
        add_textbox(s, Inches(4.3), y, Inches(0.5), Inches(0.22),
                   count, font_size=10, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
        y += Inches(0.42)

    # Right: Admin API
    add_textbox(s, Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.35),
                "ADMIN API (70 endpoints)", font_size=14, bold=True, color=ACCENT_ORANGE)

    admin_groups = [
        ("Agents", "list, get, update, delete, reload", "5"),
        ("Collections", "CRUD, documents, upload, sync, repair, health", "12"),
        ("Prompts", "list, get, update, delete", "4"),
        ("Graphs", "CRUD, build, refresh, validate, prompts, skills, tune", "12"),
        ("Access Control", "roles, permissions, bindings, memberships, principals", "14"),
        ("Config", "schema, effective, validate, apply", "4"),
        ("Uploads", "list, get, delete, reindex", "4"),
        ("MCP/System", "connections, operations, overview, architecture", "15"),
    ]

    y = Inches(1.8)
    for name, desc, count in admin_groups:
        add_shape(s, MSO_SHAPE.RECTANGLE, Inches(5.2), y, Inches(0.06), Inches(0.4), fill_color=ACCENT_ORANGE)
        add_textbox(s, Inches(5.35), y, Inches(1.3), Inches(0.22),
                   name, font_size=10, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, Inches(5.35), y + Inches(0.2), Inches(3.5), Inches(0.2),
                   desc, font_size=8, color=RGBColor(0x77, 0x77, 0x77))
        add_textbox(s, Inches(9.0), y, Inches(0.5), Inches(0.22),
                   count, font_size=10, bold=True, color=ACCENT_ORANGE, align=PP_ALIGN.RIGHT)
        y += Inches(0.42)

build_content_slide(slide, "API Surface — 117 Endpoints", body_api)


# ─── INTRO SLIDE 10: Document Processing Pipeline ──────────────────────
slide = prs.slides.add_slide(layout)
new_slides.append(slide)

def body_docs(s):
    # Left: Document subsystem
    add_textbox(s, MARGIN, Inches(0.9), Inches(4.3), Inches(0.35),
                "DOCUMENT PROCESSING (documents/)", font_size=14, bold=True, color=NAVY)
    add_textbox(s, MARGIN, Inches(1.25), Inches(4.3), Inches(0.5),
                "A dedicated subsystem for advanced document analysis beyond simple RAG retrieval.",
                font_size=11, color=RGBColor(0x66, 0x66, 0x66))

    doc_features = [
        ("Document Comparison", "Side-by-side clause-level diff between documents", TEAL),
        ("Consolidation Campaigns", "Multi-doc merge with conflict detection and resolution", ACCENT_ORANGE),
        ("Template Transform", "Map document content into structured output templates", ACCENT_PURPLE),
        ("Evidence Binder", "Compile graded evidence into formatted reference packets", ACCENT_GREEN),
        ("Requirement Extraction", "Extract and export structured requirements from specs", DARK_TEAL),
        ("Similarity Analysis", "Cross-document similarity scoring and clustering", RGBColor(0x5C, 0x6B, 0xC0)),
    ]

    y = Inches(1.85)
    for title, desc, color in doc_features:
        add_shape(s, MSO_SHAPE.RECTANGLE, MARGIN, y, Inches(0.06), Inches(0.5), fill_color=color)
        add_textbox(s, MARGIN + Inches(0.15), y, Inches(4.0), Inches(0.22),
                   title, font_size=11, bold=True, color=RGBColor(0x33, 0x33, 0x33))
        add_textbox(s, MARGIN + Inches(0.15), y + Inches(0.22), Inches(4.0), Inches(0.3),
                   desc, font_size=9, color=RGBColor(0x66, 0x66, 0x66))
        y += Inches(0.52)

    # Right: Storage + new modules
    add_textbox(s, Inches(5.2), Inches(0.9), Inches(4.3), Inches(0.35),
                "NEW MODULES SINCE V2", font_size=14, bold=True, color=NAVY)

    modules = [
        ("documents/", "7 files", "Compare, consolidate, extract, similarity, serializers, models"),
        ("mcp/", "4 files", "MCP service, client, security, registry integration"),
        ("graph/", "10 files", "Phased build, prompt tuning, community recovery, structured search"),
        ("storage/", "2 files", "Blob store abstraction layer for file artifacts"),
        ("tools/discovery.py", "1 file", "Deferred tool discovery and late-binding system"),
        ("tools/document_tools.py", "1 file", "evidence_binder, template_transform, consolidation"),
    ]

    y = Inches(1.4)
    for name, count, desc in modules:
        add_shape(s, MSO_SHAPE.RECTANGLE, Inches(5.2), y, Inches(4.3), Inches(0.65), fill_color=CARD_BG)
        add_textbox(s, Inches(5.35), y + Inches(0.05), Inches(2.5), Inches(0.22),
                   name, font_size=10, bold=True, color=WHITE)
        add_textbox(s, Inches(8.5), y + Inches(0.05), Inches(0.9), Inches(0.22),
                   count, font_size=9, color=TEAL, align=PP_ALIGN.RIGHT)
        add_textbox(s, Inches(5.35), y + Inches(0.28), Inches(4.0), Inches(0.35),
                   desc, font_size=8, color=LIGHT_GRAY)
        y += Inches(0.7)

build_content_slide(slide, "Document Processing & New Modules", body_docs)


# ══════════════════════════════════════════════════════════════
# REORDER: Move new slides to position 2-11 (after title slide)
# ══════════════════════════════════════════════════════════════

from pptx.oxml.ns import qn

# Get presentation.xml sldIdLst
pres_elem = prs.part._element
sldIdLst = pres_elem.find(qn('p:sldIdLst'))
all_sldIds = list(sldIdLst)

# The new slides are the last 10 in the list
num_new = 10
existing_ids = all_sldIds[:-num_new]
new_ids = all_sldIds[-num_new:]

# Rebuild: title slide (first) + new slides + rest of existing
reordered = [existing_ids[0]] + new_ids + existing_ids[1:]

# Clear and rebuild
for child in list(sldIdLst):
    sldIdLst.remove(child)
for sld_id in reordered:
    sldIdLst.append(sld_id)


# ══════════════════════════════════════════════════════════════
# UPDATE EXISTING ARCHITECTURE SLIDES
# ══════════════════════════════════════════════════════════════

# ── Global text replacements across ALL existing slides ──
# These fix outdated numbers, tool names, and architecture claims

replacements = {
    # Version update
    "Agentic Chatbot v2": "Agentic Chatbot v3",
    # Tool count fixes
    "39+ tools": "54 tools",
    "39 tools": "54 tools",
    "20+ tools": "54 tools",
    # Tool group fixes
    "6 groups": "8 groups",
    "6 tool groups": "8 tool groups",
    # API endpoint fixes
    "23 Public Endpoints": "47 Public Endpoints",
    "All 23 Public": "All 47 Public",
    "23 public endpoints": "47 public endpoints",
    "28 admin endpoints": "70 admin endpoints",
    # Execution mode fixes
    "8 execution modes": "7 execution modes",
    "6 execution modes": "7 execution modes",
    # Tool name fixes (old names -> actual code names)
    "python_repl": "execute_code",
    "recall_memory": "memory_load",
    '"remember"': '"memory_save"',
    "analyze_csv": "load_dataset",
    # Tool group renames
    "Utility (3)": "Utility (3)",  # same
    "Memory (3)": "Memory (3)",     # same
    # Architecture updates
    "Specialist RAG (19)": "RAG Gateway (12) + Graph Gateway (8)",
}

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old_text, new_text in replacements.items():
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
        # Also check table cells
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            for old_text, new_text in replacements.items():
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════

prs.save(OUT)
print(f"Saved to {OUT}")
print(f"Total slides: {len(prs.slides)}")
