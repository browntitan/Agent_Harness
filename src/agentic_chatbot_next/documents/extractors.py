from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Sequence

from agentic_chatbot_next.documents.models import (
    DocumentElement,
    DocumentExtractResult,
    DocumentFigure,
    DocumentIdentity,
    DocumentSection,
    DocumentTable,
)
from agentic_chatbot_next.rag.doc_targets import resolve_indexed_docs as resolve_named_indexed_docs
from agentic_chatbot_next.rag.retrieval_scope import (
    document_source_policy_requires_repository,
    has_upload_evidence,
    repository_upload_doc_ids,
    resolve_kb_collection_id,
    resolve_search_collection_ids,
    resolve_upload_collection_id,
)
from agentic_chatbot_next.storage import blob_ref_from_record, build_blob_store


SUPPORTED_FILE_TYPES = {"pdf", "docx", "pptx", "xlsx", "xls", "txt", "md", "markdown", "csv", "tsv"}
_HEADING_NUMBER_RE = re.compile(r"^\s*(?P<num>(?:\d+\.)+\d*|\d+)\s+(?P<title>[A-Z][^\n]{2,180})")
_TEXT_SPLIT_RE = re.compile(r"\n{2,}|(?<=[.!?])\s+(?=[A-Z0-9])")
_WHITESPACE_RE = re.compile(r"\s+")


class DocumentResolutionError(ValueError):
    def __init__(self, payload: Dict[str, Any]):
        super().__init__(str(payload.get("error") or "Document could not be resolved."))
        self.payload = payload


@dataclass(frozen=True)
class DocumentSource:
    identity: DocumentIdentity
    path: Path | None = None
    record: Any | None = None
    chunks_only: bool = False


def _tenant_id(settings: object, session: object) -> str:
    return str(getattr(session, "tenant_id", getattr(settings, "default_tenant_id", "local-dev")) or "local-dev")


def _safe_text(value: Any) -> str:
    return _WHITESPACE_RE.sub(" ", str(value or "").strip())


def _file_type_from_path(path: Path | None, fallback: str = "") -> str:
    if fallback:
        return str(fallback).strip().lower().lstrip(".")
    if path is None:
        return ""
    return path.suffix.lower().lstrip(".")


def _hash_path(path: Path | None) -> str:
    if path is None or not path.exists() or not path.is_file():
        return ""
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _record_identity(record: Any, *, source_scope: str, path: Path | None = None) -> DocumentIdentity:
    fallback_file_type = _file_type_from_path(path, str(getattr(record, "file_type", "") or ""))
    return DocumentIdentity(
        doc_id=str(getattr(record, "doc_id", "") or ""),
        title=str(getattr(record, "title", "") or (path.name if path else "")),
        source_type=str(getattr(record, "source_type", "") or ""),
        source_path=str(getattr(record, "source_path", "") or (str(path) if path else "")),
        collection_id=str(getattr(record, "collection_id", "") or ""),
        file_type=fallback_file_type,
        content_hash=str(getattr(record, "content_hash", "") or _hash_path(path)),
        source_scope=source_scope,
    )


def _workspace_root(session: object) -> Path | None:
    workspace = getattr(session, "workspace", None)
    if workspace is not None:
        root = getattr(workspace, "root", None)
        return Path(root) if root else None
    root = str(getattr(session, "workspace_root", "") or "").strip()
    return Path(root) if root else None


def _workspace_source(session: object, document_ref: str) -> DocumentSource | None:
    root = _workspace_root(session)
    if root is None or not root.exists():
        return None
    requested = str(document_ref or "").strip()
    if not requested:
        return None
    candidates = [requested, Path(requested).name]
    try:
        files = {path.name: path for path in root.iterdir() if path.is_file()}
    except Exception:
        files = {}
    path = next((files.get(name) for name in candidates if files.get(name)), None)
    if path is None:
        return None
    identity = DocumentIdentity(
        doc_id=f"workspace:{path.name}",
        title=path.name,
        source_type="workspace",
        source_path=str(path),
        file_type=_file_type_from_path(path),
        content_hash=_hash_path(path),
        source_scope="workspace",
    )
    return DocumentSource(identity=identity, path=path)


def _source_type_for_scope(source_scope: str) -> str:
    normalized = str(source_scope or "").strip().lower()
    if normalized == "uploads":
        return "upload"
    if normalized == "kb":
        return "kb"
    return ""


def _collection_ids_for_scope(settings: object, session: object, source_scope: str, collection_id: str) -> tuple[str, ...]:
    requested = str(collection_id or "").strip()
    normalized = str(source_scope or "").strip().lower()
    if requested:
        return (requested,)
    if normalized == "uploads":
        return (resolve_upload_collection_id(settings, session),)
    if normalized == "kb":
        return (resolve_kb_collection_id(settings, session),)
    return tuple(resolve_search_collection_ids(settings, session))


def _record_matches_scope(record: Any, *, source_scope: str, collection_id: str) -> bool:
    source_type = _source_type_for_scope(source_scope)
    if source_type and str(getattr(record, "source_type", "") or "").strip().lower() != source_type:
        return False
    if collection_id and str(getattr(record, "collection_id", "") or "") != collection_id:
        return False
    return True


def _element_location(element: DocumentElement) -> Dict[str, Any]:
    return {
        key: value
        for key, value in {
            "section_title": element.section_title,
            "clause_number": element.clause_number,
            "page_number": element.page_number,
            "slide_number": element.slide_number,
            "sheet_name": element.sheet_name,
            "row_start": element.row_start,
            "row_end": element.row_end,
            "cell_range": element.cell_range,
        }.items()
        if value not in ("", None)
    }


class _ExtractionBuilder:
    def __init__(self, document: DocumentIdentity, *, max_elements: int) -> None:
        self.document = document
        self.max_elements = max(1, int(max_elements or 200))
        self.metadata: Dict[str, Any] = {}
        self.sections: List[DocumentSection] = []
        self.elements: List[DocumentElement] = []
        self.tables: List[DocumentTable] = []
        self.figures: List[DocumentFigure] = []
        self.warnings: List[str] = []
        self._section_stack: List[DocumentSection] = []
        self._section_keys: set[str] = set()
        self.truncated = False

    def add_warning(self, text: str) -> None:
        clean = _safe_text(text)
        if clean and clean not in self.warnings:
            self.warnings.append(clean)

    def add_section(
        self,
        title: str,
        *,
        level: int = 1,
        location: Dict[str, Any] | None = None,
    ) -> DocumentSection:
        clean = _safe_text(title)
        if not clean:
            clean = "Untitled Section"
        level = max(1, int(level or 1))
        while self._section_stack and self._section_stack[-1].level >= level:
            self._section_stack.pop()
        parent_id = self._section_stack[-1].section_id if self._section_stack else ""
        key = f"{parent_id}|{level}|{clean.casefold()}"
        if key in self._section_keys:
            section = next(item for item in self.sections if f"{item.parent_id}|{item.level}|{item.title.casefold()}" == key)
        else:
            section = DocumentSection(
                section_id=f"sec_{len(self.sections) + 1:04d}",
                title=clean,
                level=level,
                parent_id=parent_id,
                order=len(self.sections),
                location=dict(location or {}),
            )
            self.sections.append(section)
            self._section_keys.add(key)
        self._section_stack.append(section)
        return section

    def current_section(self) -> DocumentSection | None:
        return self._section_stack[-1] if self._section_stack else None

    def add_element(
        self,
        element_type: str,
        text: str,
        *,
        section: DocumentSection | None = None,
        clause_number: str = "",
        page_number: int | None = None,
        slide_number: int | None = None,
        sheet_name: str = "",
        row_start: int | None = None,
        row_end: int | None = None,
        cell_range: str = "",
        metadata: Dict[str, Any] | None = None,
    ) -> None:
        clean = _safe_text(text)
        if not clean:
            return
        if len(self.elements) >= self.max_elements:
            self.truncated = True
            return
        active = section or self.current_section()
        path = self._section_path(active)
        self.elements.append(
            DocumentElement(
                element_id=f"el_{len(self.elements) + 1:05d}",
                element_type=str(element_type or "paragraph"),
                text=clean,
                order=len(self.elements),
                section_id=active.section_id if active else "",
                section_title=active.title if active else "",
                section_path=path,
                clause_number=str(clause_number or ""),
                page_number=page_number,
                slide_number=slide_number,
                sheet_name=str(sheet_name or ""),
                row_start=row_start,
                row_end=row_end,
                cell_range=str(cell_range or ""),
                metadata=dict(metadata or {}),
            )
        )

    def _section_path(self, section: DocumentSection | None) -> List[str]:
        if section is None:
            return []
        by_id = {item.section_id: item for item in self.sections}
        path: List[str] = []
        cursor: DocumentSection | None = section
        while cursor is not None:
            path.append(cursor.title)
            cursor = by_id.get(cursor.parent_id) if cursor.parent_id else None
        return list(reversed(path))

    def result(self) -> DocumentExtractResult:
        if self.truncated:
            self.add_warning(f"Element extraction was truncated at {self.max_elements} elements.")
        return DocumentExtractResult(
            document=self.document,
            metadata=self.metadata,
            sections=self.sections,
            elements=self.elements,
            tables=self.tables,
            figures=self.figures,
            warnings=self.warnings,
            truncated=self.truncated,
        )


class DocumentExtractionService:
    def __init__(self, settings: object, stores: object, session: object) -> None:
        self.settings = settings
        self.stores = stores
        self.session = session
        self.tenant_id = _tenant_id(settings, session)

    def extract(
        self,
        *,
        document_ref: str,
        source_scope: str = "auto",
        collection_id: str = "",
        include_tables: bool = True,
        include_figures: bool = True,
        include_metadata: bool = True,
        include_hierarchy: bool = True,
        max_elements: int = 200,
    ) -> DocumentExtractResult:
        source = self.resolve_source(document_ref, source_scope=source_scope, collection_id=collection_id)
        file_type = source.identity.file_type.lower()
        if source.path is not None and source.path.exists() and file_type in SUPPORTED_FILE_TYPES:
            result = self._extract_path(
                source,
                include_tables=include_tables,
                include_figures=include_figures,
                include_metadata=include_metadata,
                include_hierarchy=include_hierarchy,
                max_elements=max_elements,
            )
        else:
            result = self._extract_chunks(
                source,
                include_tables=include_tables,
                include_figures=include_figures,
                include_metadata=include_metadata,
                include_hierarchy=include_hierarchy,
                max_elements=max_elements,
            )
        return result

    def resolve_source(self, document_ref: str, *, source_scope: str = "auto", collection_id: str = "") -> DocumentSource:
        normalized_scope = str(source_scope or "auto").strip().lower()
        if normalized_scope not in {"auto", "uploads", "kb", "workspace"}:
            normalized_scope = "auto"

        if normalized_scope in {"auto", "workspace"}:
            workspace_source = _workspace_source(self.session, document_ref)
            if workspace_source is not None:
                return workspace_source
            if normalized_scope == "workspace":
                raise DocumentResolutionError(
                    {
                        "error": "Workspace document was not found.",
                        "document_ref": str(document_ref or ""),
                        "source_scope": "workspace",
                    }
                )

        if normalized_scope == "auto":
            candidate_scopes = ["uploads", "kb"] if has_upload_evidence(self.session) else ["kb", "uploads"]
            errors: list[Dict[str, Any]] = []
            for candidate_scope in candidate_scopes:
                try:
                    return self._resolve_indexed_source(
                        document_ref,
                        source_scope=candidate_scope,
                        collection_id=collection_id,
                    )
                except DocumentResolutionError as exc:
                    errors.append(exc.payload)
            raise DocumentResolutionError(
                {
                    "error": "Document reference could not be resolved in uploads or knowledge-base scopes.",
                    "document_ref": str(document_ref or ""),
                    "source_scope": "auto",
                    "attempts": errors,
                }
            )
        return self._resolve_indexed_source(
            document_ref,
            source_scope=normalized_scope,
            collection_id=collection_id,
        )

    def _resolve_indexed_source(self, document_ref: str, *, source_scope: str, collection_id: str = "") -> DocumentSource:
        effective_scope = source_scope
        collection_ids = _collection_ids_for_scope(self.settings, self.session, effective_scope, collection_id)
        source_type = _source_type_for_scope(effective_scope)
        record = self._record_by_doc_id(document_ref, source_scope=effective_scope, collection_id=str(collection_id or ""))
        if record is None:
            resolution = resolve_named_indexed_docs(
                self.stores,
                settings=self.settings,
                tenant_id=self.tenant_id,
                names=[document_ref],
                collection_ids=collection_ids,
            )
            if resolution.ambiguous or resolution.missing or not resolution.resolved:
                raise DocumentResolutionError(
                    {
                        "error": "Document reference could not be resolved.",
                        "document_ref": str(document_ref or ""),
                        "source_scope": effective_scope,
                        "collection_id": str(collection_id or ""),
                        "resolution": resolution.to_dict(),
                    }
                )
            record = self.stores.doc_store.get_document(resolution.resolved[0].doc_id, self.tenant_id)
        if record is None:
            raise DocumentResolutionError(
                {
                    "error": "Indexed document was not found.",
                    "document_ref": str(document_ref or ""),
                    "source_scope": effective_scope,
                }
            )
        if source_type and str(getattr(record, "source_type", "") or "").strip().lower() != source_type:
            raise DocumentResolutionError(
                {
                    "error": "Resolved document exists but is outside the requested source scope.",
                    "document_ref": str(document_ref or ""),
                    "requested_source_scope": effective_scope,
                    "document_source_type": str(getattr(record, "source_type", "") or ""),
                }
            )
        path_text = str(getattr(record, "source_path", "") or "")
        path = Path(path_text).expanduser() if path_text else None
        if path is None or not path.exists():
            blob_ref = blob_ref_from_record(record)
            if blob_ref is not None:
                try:
                    path = build_blob_store(self.settings).materialize_to_path(blob_ref)
                except Exception:
                    path = None
        identity = _record_identity(record, source_scope=effective_scope, path=path)
        identity.parser_path = "source_file" if path is not None and path.exists() else "indexed_chunks"
        return DocumentSource(identity=identity, path=path if path is not None and path.exists() else None, record=record, chunks_only=path is None or not path.exists())

    def _record_by_doc_id(self, document_ref: str, *, source_scope: str, collection_id: str) -> Any | None:
        doc_id = str(document_ref or "").strip()
        if not doc_id:
            return None
        try:
            record = self.stores.doc_store.get_document(doc_id, self.tenant_id)
        except Exception:
            record = None
        if record is None:
            return None
        if not _record_matches_scope(record, source_scope=source_scope, collection_id=collection_id):
            return None
        if str(getattr(record, "source_type", "") or "").strip().lower() == "upload":
            uploaded_doc_ids = set(repository_upload_doc_ids(self.session))
            if document_source_policy_requires_repository(self.session) and uploaded_doc_ids and doc_id not in uploaded_doc_ids:
                return None
        return record

    def _extract_path(
        self,
        source: DocumentSource,
        *,
        include_tables: bool,
        include_figures: bool,
        include_metadata: bool,
        include_hierarchy: bool,
        max_elements: int,
    ) -> DocumentExtractResult:
        path = source.path
        if path is None:
            return self._extract_chunks(
                source,
                include_tables=include_tables,
                include_figures=include_figures,
                include_metadata=include_metadata,
                include_hierarchy=include_hierarchy,
                max_elements=max_elements,
            )
        file_type = source.identity.file_type.lower()
        if file_type == "pdf":
            return self._extract_pdf(source, include_figures=include_figures, include_metadata=include_metadata, max_elements=max_elements)
        if file_type == "docx":
            return self._extract_docx(source, include_tables=include_tables, include_figures=include_figures, include_metadata=include_metadata, max_elements=max_elements)
        if file_type == "pptx":
            return self._extract_pptx(source, include_tables=include_tables, include_figures=include_figures, include_metadata=include_metadata, max_elements=max_elements)
        if file_type in {"xlsx", "xls"}:
            return self._extract_workbook(source, include_tables=include_tables, include_metadata=include_metadata, max_elements=max_elements)
        return self._extract_text_file(source, include_hierarchy=include_hierarchy, include_metadata=include_metadata, max_elements=max_elements)

    def _extract_chunks(
        self,
        source: DocumentSource,
        *,
        include_tables: bool,
        include_figures: bool,
        include_metadata: bool,
        include_hierarchy: bool,
        max_elements: int,
    ) -> DocumentExtractResult:
        builder = _ExtractionBuilder(source.identity, max_elements=max_elements)
        builder.document.parser_path = "indexed_chunks"
        if include_metadata:
            builder.metadata.update(
                {
                    "doc_id": source.identity.doc_id,
                    "title": source.identity.title,
                    "content_hash": source.identity.content_hash,
                    "num_chunks": int(getattr(source.record, "num_chunks", 0) or 0),
                }
            )
        try:
            chunks = list(self.stores.chunk_store.list_document_chunks(source.identity.doc_id, self.tenant_id))
        except Exception:
            chunks = []
        if not chunks:
            builder.add_warning("No source file or indexed chunks were available for structured extraction.")
            return builder.result()
        section_by_title: Dict[str, DocumentSection] = {}
        sheet_rows: Dict[str, List[List[str]]] = {}
        for chunk in sorted(chunks, key=lambda item: int(getattr(item, "chunk_index", 0) or 0)):
            section = None
            section_title = str(getattr(chunk, "section_title", "") or "")
            if include_hierarchy and section_title:
                section = section_by_title.get(section_title)
                if section is None:
                    section = builder.add_section(section_title, level=1, location={"chunk_index": getattr(chunk, "chunk_index", None)})
                    section_by_title[section_title] = section
            sheet_name = str(getattr(chunk, "sheet_name", "") or "")
            content = str(getattr(chunk, "content", "") or "")
            builder.add_element(
                str(getattr(chunk, "chunk_type", "") or "chunk"),
                content,
                section=section,
                clause_number=str(getattr(chunk, "clause_number", "") or ""),
                page_number=getattr(chunk, "page_number", None),
                sheet_name=sheet_name,
                row_start=getattr(chunk, "row_start", None),
                row_end=getattr(chunk, "row_end", None),
                cell_range=str(getattr(chunk, "cell_range", "") or ""),
                metadata={"chunk_id": str(getattr(chunk, "chunk_id", "") or ""), "chunk_index": int(getattr(chunk, "chunk_index", 0) or 0)},
            )
            if include_tables and sheet_name:
                sheet_rows.setdefault(sheet_name, []).append([content])
        if include_tables:
            for sheet_name, rows in sheet_rows.items():
                builder.tables.append(
                    DocumentTable(
                        table_id=f"table_{len(builder.tables) + 1:04d}",
                        title=sheet_name,
                        sheet_name=sheet_name,
                        rows=rows[:100],
                        metadata={"source": "indexed_chunks", "row_count": len(rows)},
                    )
                )
        if include_figures:
            del include_figures
        builder.add_warning("Used indexed text chunks because the original source file was unavailable.")
        return builder.result()

    def _extract_pdf(
        self,
        source: DocumentSource,
        *,
        include_figures: bool,
        include_metadata: bool,
        max_elements: int,
    ) -> DocumentExtractResult:
        path = source.path
        assert path is not None
        docling_result = self._try_docling(source, max_elements=max_elements) if bool(getattr(self.settings, "docling_enabled", False)) else None
        if docling_result is not None:
            return docling_result

        builder = _ExtractionBuilder(source.identity, max_elements=max_elements)
        builder.document.parser_path = "pypdf_text"
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            if include_metadata and getattr(reader, "metadata", None):
                builder.metadata.update({str(key).lstrip("/"): str(value) for key, value in dict(reader.metadata or {}).items()})
            builder.metadata["page_count"] = len(reader.pages)
            for index, page in enumerate(reader.pages, start=1):
                page_section = builder.add_section(f"Page {index}", level=1, location={"page_number": index})
                text = page.extract_text() or ""
                for part in self._text_blocks(text):
                    heading = _HEADING_NUMBER_RE.match(part)
                    if heading:
                        section = builder.add_section(heading.group("title"), level=max(1, heading.group("num").count(".") + 1), location={"page_number": index})
                        builder.add_element("heading", part, section=section, clause_number=heading.group("num"), page_number=index)
                    else:
                        builder.add_element("paragraph", part, section=page_section, page_number=index)
        except Exception as exc:
            builder.add_warning(f"PDF text extraction failed: {exc}")
        if include_figures:
            self._add_pdf_figures(path, builder)
        builder.add_warning("text_only_fallback: PDF layout/table extraction was not available without Docling.")
        return builder.result()

    def _try_docling(self, source: DocumentSource, *, max_elements: int) -> DocumentExtractResult | None:
        path = source.path
        if path is None:
            return None
        builder = _ExtractionBuilder(source.identity, max_elements=max_elements)
        try:
            from docling.document_converter import DocumentConverter

            converted = DocumentConverter().convert(str(path))
            document = converted.document
            markdown = document.export_to_markdown()
            builder.document.parser_path = "docling"
            builder.metadata["docling_available"] = True
            for part in markdown.splitlines():
                clean = part.strip()
                if not clean:
                    continue
                if clean.startswith("#"):
                    level = max(1, len(clean) - len(clean.lstrip("#")))
                    section = builder.add_section(clean.lstrip("#").strip(), level=level)
                    builder.add_element("heading", section.title, section=section)
                else:
                    builder.add_element("paragraph", clean)
            return builder.result()
        except Exception as exc:
            builder.add_warning(f"Docling extraction failed; falling back to native parser: {exc}")
            return None

    def _add_pdf_figures(self, path: Path, builder: _ExtractionBuilder) -> None:
        try:
            import fitz

            with fitz.open(path) as document:
                for page_index in range(len(document)):
                    images = document[page_index].get_images(full=True)
                    for image_index, image in enumerate(images, start=1):
                        builder.figures.append(
                            DocumentFigure(
                                figure_id=f"fig_{len(builder.figures) + 1:04d}",
                                title=f"Page {page_index + 1} image {image_index}",
                                page_number=page_index + 1,
                                description="Embedded PDF image.",
                                metadata={"xref": image[0] if image else None},
                            )
                        )
        except Exception:
            return

    def _extract_docx(
        self,
        source: DocumentSource,
        *,
        include_tables: bool,
        include_figures: bool,
        include_metadata: bool,
        max_elements: int,
    ) -> DocumentExtractResult:
        path = source.path
        assert path is not None
        builder = _ExtractionBuilder(source.identity, max_elements=max_elements)
        builder.document.parser_path = "python_docx"
        try:
            from docx import Document

            document = Document(str(path))
            if include_metadata:
                builder.metadata.update(self._docx_core_properties(document.core_properties))
            for paragraph in document.paragraphs:
                text = paragraph.text or ""
                style_name = str(getattr(getattr(paragraph, "style", None), "name", "") or "")
                level = self._heading_level(style_name)
                if level:
                    section = builder.add_section(text, level=level)
                    builder.add_element("heading", text, section=section)
                else:
                    builder.add_element("paragraph", text)
            if include_tables:
                for index, table in enumerate(document.tables, start=1):
                    rows = [[_safe_text(cell.text) for cell in row.cells] for row in table.rows]
                    columns = rows[0] if rows else []
                    body_rows = rows[1:] if len(rows) > 1 else rows
                    builder.tables.append(
                        DocumentTable(
                            table_id=f"table_{index:04d}",
                            title=f"Table {index}",
                            columns=columns,
                            rows=body_rows[:100],
                            metadata={"row_count": len(body_rows), "parser": "python_docx"},
                        )
                    )
                    for row_index, row in enumerate(body_rows, start=1):
                        row_text = "; ".join(item for item in row if item)
                        builder.add_element("table_row", row_text, metadata={"table_id": f"table_{index:04d}", "row_index": row_index})
            if include_figures:
                for index, shape in enumerate(document.inline_shapes, start=1):
                    builder.figures.append(
                        DocumentFigure(
                            figure_id=f"fig_{index:04d}",
                            title=f"Inline figure {index}",
                            description="DOCX inline shape.",
                            metadata={"height": int(getattr(shape, "height", 0) or 0), "width": int(getattr(shape, "width", 0) or 0)},
                        )
                    )
        except Exception as exc:
            builder.add_warning(f"DOCX extraction failed: {exc}")
        return builder.result()

    def _docx_core_properties(self, props: Any) -> Dict[str, Any]:
        fields = ["author", "category", "comments", "content_status", "created", "identifier", "keywords", "language", "last_modified_by", "modified", "revision", "subject", "title", "version"]
        payload: Dict[str, Any] = {}
        for field in fields:
            value = getattr(props, field, None)
            if value not in ("", None):
                payload[field] = str(value)
        return payload

    def _heading_level(self, style_name: str) -> int:
        match = re.search(r"heading\s+(\d+)", str(style_name or ""), re.IGNORECASE)
        return int(match.group(1)) if match else 0

    def _extract_pptx(
        self,
        source: DocumentSource,
        *,
        include_tables: bool,
        include_figures: bool,
        include_metadata: bool,
        max_elements: int,
    ) -> DocumentExtractResult:
        path = source.path
        assert path is not None
        builder = _ExtractionBuilder(source.identity, max_elements=max_elements)
        builder.document.parser_path = "python_pptx"
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            presentation = Presentation(str(path))
            if include_metadata:
                builder.metadata.update(self._pptx_core_properties(presentation.core_properties))
                builder.metadata["slide_count"] = len(presentation.slides)
            for slide_index, slide in enumerate(presentation.slides, start=1):
                section = builder.add_section(f"Slide {slide_index}", level=1, location={"slide_number": slide_index})
                for shape_index, shape in enumerate(slide.shapes, start=1):
                    if getattr(shape, "has_text_frame", False):
                        text = _safe_text(getattr(shape, "text", ""))
                        if text:
                            element_type = "slide_title" if shape_index == 1 else "slide_text"
                            builder.add_element(element_type, text, section=section, slide_number=slide_index, metadata={"shape_index": shape_index})
                    if include_tables and getattr(shape, "has_table", False):
                        rows = [[_safe_text(cell.text) for cell in row.cells] for row in shape.table.rows]
                        columns = rows[0] if rows else []
                        body_rows = rows[1:] if len(rows) > 1 else rows
                        table_id = f"table_{len(builder.tables) + 1:04d}"
                        builder.tables.append(
                            DocumentTable(
                                table_id=table_id,
                                title=f"Slide {slide_index} table",
                                slide_number=slide_index,
                                columns=columns,
                                rows=body_rows[:100],
                                metadata={"row_count": len(body_rows), "shape_index": shape_index},
                            )
                        )
                        for row_index, row in enumerate(body_rows, start=1):
                            builder.add_element("table_row", "; ".join(item for item in row if item), section=section, slide_number=slide_index, metadata={"table_id": table_id, "row_index": row_index})
                    if include_figures and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                        builder.figures.append(
                            DocumentFigure(
                                figure_id=f"fig_{len(builder.figures) + 1:04d}",
                                title=f"Slide {slide_index} picture",
                                slide_number=slide_index,
                                description="PPTX picture shape.",
                                metadata={"shape_index": shape_index},
                            )
                        )
                notes = self._slide_notes_text(slide)
                if notes:
                    builder.add_element("speaker_notes", notes, section=section, slide_number=slide_index)
        except Exception as exc:
            builder.add_warning(f"PPTX extraction failed: {exc}")
        return builder.result()

    def _pptx_core_properties(self, props: Any) -> Dict[str, Any]:
        fields = ["author", "category", "comments", "content_status", "created", "identifier", "keywords", "language", "last_modified_by", "modified", "revision", "subject", "title", "version"]
        payload: Dict[str, Any] = {}
        for field in fields:
            value = getattr(props, field, None)
            if value not in ("", None):
                payload[field] = str(value)
        return payload

    def _slide_notes_text(self, slide: Any) -> str:
        try:
            if not getattr(slide, "has_notes_slide", False):
                return ""
            return _safe_text(getattr(slide.notes_slide.notes_text_frame, "text", "") or "")
        except Exception:
            return ""

    def _extract_workbook(
        self,
        source: DocumentSource,
        *,
        include_tables: bool,
        include_metadata: bool,
        max_elements: int,
    ) -> DocumentExtractResult:
        path = source.path
        assert path is not None
        if path.suffix.lower() == ".xls":
            return self._extract_xls_with_pandas(source, include_tables=include_tables, include_metadata=include_metadata, max_elements=max_elements)
        builder = _ExtractionBuilder(source.identity, max_elements=max_elements)
        builder.document.parser_path = "openpyxl"
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(str(path), data_only=False, read_only=False)
            try:
                if include_metadata:
                    props = workbook.properties
                    builder.metadata.update({key: str(value) for key, value in vars(props).items() if value not in ("", None) and not key.startswith("_")})
                    builder.metadata["sheet_count"] = len(workbook.worksheets)
                for sheet in workbook.worksheets:
                    section = builder.add_section(sheet.title, level=1, location={"sheet_name": sheet.title})
                    rows = list(sheet.iter_rows(values_only=False))
                    values = [[_safe_text(cell.value) for cell in row] for row in rows]
                    nonempty_rows = [(index + 1, row) for index, row in enumerate(values) if any(row)]
                    cell_range = sheet.calculate_dimension()
                    if include_tables:
                        table_id = f"table_{len(builder.tables) + 1:04d}"
                        table_columns = nonempty_rows[0][1] if nonempty_rows else []
                        table_rows = [row for _, row in nonempty_rows[1:]]
                        builder.tables.append(
                            DocumentTable(
                                table_id=table_id,
                                title=sheet.title,
                                sheet_name=sheet.title,
                                cell_range=cell_range,
                                columns=table_columns,
                                rows=table_rows[:100],
                                metadata={
                                    "row_count": len(table_rows),
                                    "max_row": sheet.max_row,
                                    "max_column": sheet.max_column,
                                    "merged_cells": [str(item) for item in sheet.merged_cells.ranges],
                                    "defined_tables": list(getattr(sheet, "tables", {}).keys()),
                                },
                            )
                        )
                    for row_number, row in nonempty_rows:
                        row_text = "; ".join(cell for cell in row if cell)
                        end_col = self._excel_column_label(max(0, len(row) - 1))
                        builder.add_element(
                            "spreadsheet_row",
                            row_text,
                            section=section,
                            sheet_name=sheet.title,
                            row_start=row_number,
                            row_end=row_number,
                            cell_range=f"A{row_number}:{end_col}{row_number}",
                        )
            finally:
                workbook.close()
        except Exception as exc:
            builder.add_warning(f"Workbook extraction failed: {exc}")
        return builder.result()

    def _extract_xls_with_pandas(
        self,
        source: DocumentSource,
        *,
        include_tables: bool,
        include_metadata: bool,
        max_elements: int,
    ) -> DocumentExtractResult:
        path = source.path
        assert path is not None
        builder = _ExtractionBuilder(source.identity, max_elements=max_elements)
        builder.document.parser_path = "pandas_xls"
        try:
            import pandas as pd

            workbook = pd.ExcelFile(path)
            if include_metadata:
                builder.metadata["sheet_count"] = len(workbook.sheet_names)
            for sheet_name in workbook.sheet_names:
                frame = workbook.parse(sheet_name=sheet_name, header=None)
                section = builder.add_section(sheet_name, level=1, location={"sheet_name": sheet_name})
                rows = [[_safe_text(value) for value in values] for values in frame.itertuples(index=False, name=None)]
                if include_tables:
                    columns = rows[0] if rows else []
                    body_rows = rows[1:] if len(rows) > 1 else rows
                    builder.tables.append(
                        DocumentTable(
                            table_id=f"table_{len(builder.tables) + 1:04d}",
                            title=sheet_name,
                            sheet_name=sheet_name,
                            columns=columns,
                            rows=body_rows[:100],
                            metadata={"row_count": len(body_rows), "parser": "pandas"},
                        )
                    )
                for row_index, row in enumerate(rows, start=1):
                    builder.add_element("spreadsheet_row", "; ".join(cell for cell in row if cell), section=section, sheet_name=sheet_name, row_start=row_index, row_end=row_index)
        except Exception as exc:
            builder.add_warning(f"XLS extraction failed: {exc}")
        return builder.result()

    def _excel_column_label(self, index: int) -> str:
        index = max(0, int(index))
        letters = ""
        while True:
            index, remainder = divmod(index, 26)
            letters = chr(ord("A") + remainder) + letters
            if index == 0:
                break
            index -= 1
        return letters

    def _extract_text_file(
        self,
        source: DocumentSource,
        *,
        include_hierarchy: bool,
        include_metadata: bool,
        max_elements: int,
    ) -> DocumentExtractResult:
        path = source.path
        assert path is not None
        builder = _ExtractionBuilder(source.identity, max_elements=max_elements)
        builder.document.parser_path = "plain_text"
        if include_metadata:
            builder.metadata["size_bytes"] = path.stat().st_size
        text = path.read_text(encoding="utf-8", errors="replace")
        for part in self._text_blocks(text):
            heading = _HEADING_NUMBER_RE.match(part)
            if include_hierarchy and heading:
                section = builder.add_section(heading.group("title"), level=max(1, heading.group("num").count(".") + 1))
                builder.add_element("heading", part, section=section, clause_number=heading.group("num"))
            elif include_hierarchy and part.startswith("#"):
                level = max(1, len(part) - len(part.lstrip("#")))
                section = builder.add_section(part.lstrip("#").strip(), level=level)
                builder.add_element("heading", section.title, section=section)
            else:
                builder.add_element("paragraph", part)
        return builder.result()

    def _text_blocks(self, text: str) -> List[str]:
        blocks = [_safe_text(part) for part in _TEXT_SPLIT_RE.split(str(text or ""))]
        return [part for part in blocks if part]


def element_location(element: DocumentElement) -> Dict[str, Any]:
    return _element_location(element)


__all__ = [
    "DocumentExtractionService",
    "DocumentResolutionError",
    "DocumentSource",
    "SUPPORTED_FILE_TYPES",
    "element_location",
]
