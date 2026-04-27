from __future__ import annotations

import csv
import io
import json
import re
from pathlib import Path
from typing import Any, Dict, List

from agentic_chatbot_next.documents.templates.models import TemplateSlide, TemplateTransformResult
from agentic_chatbot_next.runtime.artifacts import register_workspace_artifact


_FILENAME_SAFE_RE = re.compile(r"[^A-Za-z0-9_.-]+")


def safe_stem(value: str, *, fallback: str = "template_transform") -> str:
    stem = Path(str(value or "")).stem or fallback
    clean = _FILENAME_SAFE_RE.sub("_", stem).strip("._")
    return (clean or fallback)[:72]


def ensure_workspace(session: object) -> object:
    workspace = getattr(session, "workspace", None)
    if workspace is not None:
        return workspace
    workspace_root = str(getattr(session, "workspace_root", "") or "").strip()
    session_id = str(getattr(session, "session_id", "") or "").strip()
    if not workspace_root or not session_id:
        raise ValueError("No session workspace is available.")
    from agentic_chatbot_next.sandbox.workspace import SessionWorkspace

    workspace = SessionWorkspace(session_id=session_id, root=Path(workspace_root))
    workspace.open()
    session.workspace = workspace
    return workspace


def _safe_filename(filename: str) -> str:
    clean = str(filename or "").strip()
    if not clean:
        raise ValueError("filename must not be empty")
    if "/" in clean or "\\" in clean or "\x00" in clean or clean.startswith(".."):
        raise ValueError(f"unsafe artifact filename: {filename!r}")
    return clean[:128]


def write_text_artifact(session: object, filename: str, content: str) -> Dict[str, Any]:
    workspace = ensure_workspace(session)
    clean = _safe_filename(filename)
    workspace.write_text(clean, content)
    return register_workspace_artifact(session, filename=clean, label=clean)


def write_binary_artifact(session: object, filename: str, content: bytes) -> Dict[str, Any]:
    workspace = ensure_workspace(session)
    clean = _safe_filename(filename)
    path = workspace.root / clean
    path.write_bytes(content)
    return register_workspace_artifact(session, filename=clean, label=clean)


def markdown_for_result(result: TemplateTransformResult) -> str:
    lines = [
        f"# {result.template_type.replace('_', ' ').title()}",
        "",
        f"- Transform ID: {result.transform_id}",
        f"- Status: {result.status}",
        f"- Output format: {result.output_format}",
        f"- Source documents: {result.selected_document_count}",
    ]
    if result.warnings:
        lines.extend(["", "## Warnings"])
        lines.extend(f"- {warning}" for warning in result.warnings)
    if result.sections:
        lines.append("")
        for section in result.sections:
            lines.append(f"## {section.title}")
            if section.body:
                lines.append(section.body)
            for bullet in section.bullets:
                lines.append(f"- {bullet}")
            if section.source_ids:
                lines.append(f"_Sources: {', '.join(section.source_ids)}_")
            lines.append("")
    for table in result.tables:
        lines.extend(["", f"## {table.name}", ""])
        lines.append("| " + " | ".join(table.columns) + " |")
        lines.append("| " + " | ".join("---" for _ in table.columns) + " |")
        for row in table.rows:
            values = [str(row.get(column, "") or "").replace("|", "\\|") for column in table.columns]
            lines.append("| " + " | ".join(values) + " |")
    if result.slides:
        lines.extend(["", "## Slides"])
        for index, slide in enumerate(result.slides, start=1):
            lines.append(f"### Slide {index}: {slide.title}")
            lines.extend(f"- {bullet}" for bullet in slide.bullets)
            if slide.source_ids:
                lines.append(f"_Sources: {', '.join(slide.source_ids)}_")
    return "\n".join(lines).rstrip() + "\n"


def csv_for_table(result: TemplateTransformResult) -> str:
    table = result.tables[0] if result.tables else None
    if table is None:
        return ""
    buffer = io.StringIO()
    writer = csv.DictWriter(buffer, fieldnames=table.columns)
    writer.writeheader()
    for row in table.rows:
        writer.writerow({column: row.get(column, "") for column in table.columns})
    return buffer.getvalue()


def source_trace_json(result: TemplateTransformResult) -> str:
    payload = {
        "transform_id": result.transform_id,
        "template_type": result.template_type,
        "source_documents": [item.to_dict() for item in result.source_documents],
    }
    return json.dumps(payload, indent=2, sort_keys=True)


def sidecar_json(result: TemplateTransformResult) -> str:
    return json.dumps(result.to_dict(), indent=2, sort_keys=True)


def render_docx_bytes(result: TemplateTransformResult) -> bytes:
    from docx import Document

    document = Document()
    document.add_heading(result.template_type.replace("_", " ").title(), level=0)
    for section in result.sections:
        document.add_heading(section.title, level=1)
        if section.body:
            document.add_paragraph(section.body)
        for bullet in section.bullets:
            document.add_paragraph(bullet, style="List Bullet")
        if section.source_ids:
            document.add_paragraph("Sources: " + ", ".join(section.source_ids))
    if result.tables:
        for table in result.tables:
            document.add_heading(table.name, level=1)
            doc_table = document.add_table(rows=1, cols=max(1, len(table.columns)))
            doc_table.style = "Table Grid"
            for index, column in enumerate(table.columns):
                doc_table.rows[0].cells[index].text = column
            for row in table.rows:
                cells = doc_table.add_row().cells
                for index, column in enumerate(table.columns):
                    cells[index].text = str(row.get(column, "") or "")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def render_xlsx_bytes(result: TemplateTransformResult) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    workbook = Workbook()
    if result.tables:
        workbook.remove(workbook.active)
        for table in result.tables:
            sheet = workbook.create_sheet(title=table.name[:31] or "Transform")
            sheet.append(table.columns)
            for cell in sheet[1]:
                cell.font = Font(bold=True)
                cell.fill = PatternFill("solid", fgColor="D9EAF7")
            for row in table.rows:
                sheet.append([row.get(column, "") for column in table.columns])
            for column_cells in sheet.columns:
                max_len = max(len(str(cell.value or "")) for cell in column_cells)
                sheet.column_dimensions[column_cells[0].column_letter].width = min(max(12, max_len + 2), 48)
    else:
        sheet = workbook.active
        sheet.title = "Transform"
        sheet.append(["Section", "Text", "Sources"])
        for section in result.sections:
            sheet.append([section.title, section.body or " ".join(section.bullets), ", ".join(section.source_ids)])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def render_pptx_bytes(result: TemplateTransformResult) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    for slide_model in result.slides or _slides_from_sections(result):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_model.title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        bullets = slide_model.bullets or [slide_model.speaker_notes or ""]
        for index, bullet in enumerate(bullets[:6]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = str(bullet or "")[:260]
            paragraph.level = 0
        if slide_model.source_ids:
            notes = slide.notes_slide.notes_text_frame
            notes.text = "Sources: " + ", ".join(slide_model.source_ids)
    if not presentation.slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = "Executive Brief"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _slides_from_sections(result: TemplateTransformResult):
    return [
        TemplateSlide(
            title=section.title,
            bullets=section.bullets or [section.body],
            source_ids=section.source_ids,
        )
        for section in result.sections[:8]
    ]


__all__ = [
    "csv_for_table",
    "markdown_for_result",
    "render_docx_bytes",
    "render_pptx_bytes",
    "render_xlsx_bytes",
    "safe_stem",
    "sidecar_json",
    "source_trace_json",
    "write_binary_artifact",
    "write_text_artifact",
]
