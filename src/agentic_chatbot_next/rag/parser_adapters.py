from __future__ import annotations

import csv
import importlib.metadata
import logging
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, List

from langchain_core.documents import Document

from agentic_chatbot_next.config import Settings
from agentic_chatbot_next.rag.ocr import IMAGE_SUFFIXES, load_image_documents, load_pdf_documents_with_ocr
from agentic_chatbot_next.rag.workbook_loader import load_workbook_documents

logger = logging.getLogger(__name__)

DOCLING_SUFFIXES = {".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv"}
TEXT_SUFFIXES = {
    ".md",
    ".markdown",
    ".txt",
    ".text",
    ".json",
    ".jsonl",
    ".yaml",
    ".yml",
    ".log",
    ".rst",
}
CSV_SUFFIXES = {".csv", ".tsv"}


@dataclass(frozen=True)
class ParserProvenanceStep:
    name: str
    status: str
    version: str = ""
    detail: str = ""
    warnings: List[str] = field(default_factory=list)
    stats: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


@dataclass(frozen=True)
class ParsedDocumentBundle:
    documents: List[Document] = field(default_factory=list)
    parser_chain: List[str] = field(default_factory=list)
    parser_provenance: dict[str, Any] = field(default_factory=dict)
    warnings: List[str] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)
    ocr_page_count: int = 0
    extraction_status: str = "success"
    extraction_error: str = ""

    def to_source_metadata(self) -> dict[str, Any]:
        return {
            "parser_chain": list(self.parser_chain),
            "parser_provenance": dict(self.parser_provenance),
            "parser_warnings": list(self.warnings),
            "parser_errors": list(self.errors),
            "ocr_page_count": int(self.ocr_page_count or 0),
            "extraction_status": self.extraction_status,
            "extraction_error": self.extraction_error,
        }


def load_documents_with_parsers(
    path: Path,
    settings: Settings,
    *,
    parser_strategy: str = "docling_primary",
) -> ParsedDocumentBundle:
    """Load a file through the corpus ingestion parser adapter layer."""
    path = Path(path)
    normalized_strategy = str(parser_strategy or "docling_primary").strip().lower() or "docling_primary"
    warnings: list[str] = []
    errors: list[str] = []
    steps: list[ParserProvenanceStep] = []
    suffix = path.suffix.lower()

    if normalized_strategy == "docling_primary" and suffix in DOCLING_SUFFIXES:
        docling_docs, docling_step = _try_docling(path)
        steps.append(docling_step)
        warnings.extend(docling_step.warnings)
        if docling_docs:
            return _bundle(
                path,
                docling_docs,
                parser_strategy=normalized_strategy,
                steps=steps,
                warnings=warnings,
                errors=errors,
            )

    try:
        native_docs, native_steps, ocr_page_count = _load_native(path, settings)
        steps.extend(native_steps)
        for step in native_steps:
            warnings.extend(step.warnings)
        if native_docs:
            return _bundle(
                path,
                native_docs,
                parser_strategy=normalized_strategy,
                steps=steps,
                warnings=warnings,
                errors=errors,
                ocr_page_count=ocr_page_count,
            )
        errors.append("No parser returned extractable content.")
    except Exception as exc:
        message = f"{type(exc).__name__}: {exc}"
        errors.append(message)
        steps.append(ParserProvenanceStep(name="native", status="failed", detail=message))
        logger.warning("Native extraction failed for %s: %s", path, exc)

    return _bundle(
        path,
        [],
        parser_strategy=normalized_strategy,
        steps=steps,
        warnings=warnings,
        errors=errors,
        extraction_status="failed",
        extraction_error="; ".join(errors)[:2000],
    )


def _bundle(
    path: Path,
    docs: list[Document],
    *,
    parser_strategy: str,
    steps: list[ParserProvenanceStep],
    warnings: list[str],
    errors: list[str],
    ocr_page_count: int = 0,
    extraction_status: str = "",
    extraction_error: str = "",
) -> ParsedDocumentBundle:
    chain = _parser_chain(docs, steps)
    provenance = {
        "strategy": parser_strategy,
        "chain": chain,
        "steps": [step.to_dict() for step in steps],
        "source_suffix": path.suffix.lower(),
    }
    status = extraction_status or ("success" if docs else "failed")
    error = extraction_error or ("" if docs else "; ".join(errors)[:2000])
    for doc in docs:
        metadata = dict(doc.metadata or {})
        parser_name = str(metadata.get("parser") or (chain[-1] if chain else "unknown"))
        doc.metadata = {
            **metadata,
            "source": metadata.get("source") or str(path),
            "parser": parser_name,
            "parser_chain": chain,
            "parser_strategy": parser_strategy,
        }
    return ParsedDocumentBundle(
        documents=docs,
        parser_chain=chain,
        parser_provenance=provenance,
        warnings=_unique(warnings),
        errors=_unique(errors),
        ocr_page_count=int(ocr_page_count or 0),
        extraction_status=status,
        extraction_error=error,
    )


def _parser_chain(docs: list[Document], steps: list[ParserProvenanceStep]) -> list[str]:
    chain: list[str] = []
    for step in steps:
        if step.status == "success" and step.name not in chain:
            chain.append(step.name)
    for doc in docs:
        name = str((doc.metadata or {}).get("parser") or "").strip()
        if name and name not in chain:
            chain.append(name)
        ocr_source = str((doc.metadata or {}).get("ocr_source") or "").strip()
        if ocr_source and ocr_source not in chain:
            chain.append(ocr_source)
    return chain


def _try_docling(path: Path) -> tuple[list[Document], ParserProvenanceStep]:
    version = _package_version("docling")
    try:
        from docling.document_converter import DocumentConverter
    except Exception as exc:
        return (
            [],
            ParserProvenanceStep(
                name="docling",
                status="unavailable",
                version=version,
                detail=str(exc),
                warnings=[f"Docling unavailable; using native fallback: {exc}"],
            ),
        )

    try:
        converter = DocumentConverter()
        result = converter.convert(str(path))
        document = getattr(result, "document", None)
        if document is None:
            return (
                [],
                ParserProvenanceStep(
                    name="docling",
                    status="empty",
                    version=version,
                    warnings=["Docling returned no document object; using native fallback."],
                ),
            )

        rendered_text = ""
        for method_name in ("export_to_markdown", "export_to_text"):
            method = getattr(document, method_name, None)
            if not callable(method):
                continue
            try:
                content = method()
            except TypeError:
                continue
            if str(content or "").strip():
                rendered_text = str(content)
                break
        if not rendered_text:
            rendered_text = str(document or "")
        if not rendered_text.strip():
            return (
                [],
                ParserProvenanceStep(
                    name="docling",
                    status="empty",
                    version=version,
                    warnings=["Docling returned empty text; using native fallback."],
                ),
            )
        return (
            [Document(page_content=rendered_text, metadata={"parser": "docling"})],
            ParserProvenanceStep(
                name="docling",
                status="success",
                version=version,
                stats={"char_count": len(rendered_text)},
            ),
        )
    except Exception as exc:
        logger.warning("Docling extraction failed for %s: %s", path, exc)
        return (
            [],
            ParserProvenanceStep(
                name="docling",
                status="failed",
                version=version,
                detail=str(exc),
                warnings=[f"Docling extraction failed; using native fallback: {exc}"],
            ),
        )


def _load_native(path: Path, settings: Settings) -> tuple[list[Document], list[ParserProvenanceStep], int]:
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return _load_text(path)
    if suffix in CSV_SUFFIXES:
        return _load_csv(path)
    if suffix == ".pdf":
        return _load_pdf(path, settings)
    if suffix == ".docx":
        return _load_docx(path)
    if suffix == ".pptx":
        return _load_pptx(path)
    if suffix in {".xlsx", ".xls"}:
        return _load_workbook(path)
    if suffix in IMAGE_SUFFIXES:
        return _load_image(path, settings)
    return _load_text(path)


def _load_text(path: Path) -> tuple[list[Document], list[ParserProvenanceStep], int]:
    from langchain_community.document_loaders import TextLoader

    try:
        docs = TextLoader(str(path), encoding="utf-8", autodetect_encoding=True).load()
    except TypeError:
        docs = TextLoader(str(path), encoding="utf-8").load()
    for doc in docs:
        doc.metadata = {**(doc.metadata or {}), "parser": "text"}
    return docs, [ParserProvenanceStep(name="text", status="success", version=_package_version("langchain-community"))], 0


def _load_csv(path: Path) -> tuple[list[Document], list[ParserProvenanceStep], int]:
    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    docs: list[Document] = []
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        sample = handle.read(4096)
        handle.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
            delimiter = dialect.delimiter
        except Exception:
            pass
        reader = csv.reader(handle, delimiter=delimiter)
        rows = [row for row in reader]

    nonempty = [row for row in rows if any(str(cell or "").strip() for cell in row)]
    if not nonempty:
        return [], [ParserProvenanceStep(name="csv", status="empty", version="stdlib")], 0
    headers = [str(cell or "").strip() or f"Column {index + 1}" for index, cell in enumerate(nonempty[0])]
    summary = f"Table: {path.name}\nColumns: {', '.join(headers)}\nRows: {max(0, len(nonempty) - 1)}"
    docs.append(
        Document(
            page_content=summary,
            metadata={
                "parser": "csv",
                "chunk_index": 0,
                "chunk_type": "worksheet_summary",
                "sheet_name": path.stem,
                "row_start": 1,
                "row_end": len(nonempty),
                "cell_range": f"A1:{_column_label(max(0, len(headers) - 1))}{len(nonempty)}",
                "is_prechunked": True,
            },
        )
    )
    for index, row in enumerate(nonempty[1:], start=2):
        normalized = list(row) + [""] * max(0, len(headers) - len(row))
        pairs = [
            f"{header}: {str(value or '').strip()}"
            for header, value in zip(headers, normalized)
            if str(value or "").strip()
        ]
        if not pairs:
            continue
        docs.append(
            Document(
                page_content=f"Table: {path.name} | Row {index}: " + "; ".join(pairs[:24]),
                metadata={
                    "parser": "csv",
                    "chunk_index": len(docs),
                    "chunk_type": "spreadsheet_row",
                    "sheet_name": path.stem,
                    "row_start": index,
                    "row_end": index,
                    "cell_range": f"A{index}:{_column_label(max(0, len(headers) - 1))}{index}",
                    "is_prechunked": True,
                },
            )
        )
    return docs, [ParserProvenanceStep(name="csv", status="success", version="stdlib", stats={"row_count": len(nonempty)})], 0


def _load_pdf(path: Path, settings: Settings) -> tuple[list[Document], list[ParserProvenanceStep], int]:
    if bool(getattr(settings, "ocr_enabled", False)):
        docs = load_pdf_documents_with_ocr(
            path,
            min_page_chars=int(getattr(settings, "ocr_min_page_chars", 50) or 50),
            language=str(getattr(settings, "ocr_language", "en") or "en"),
            use_gpu=bool(getattr(settings, "ocr_use_gpu", False)),
        )
        ocr_pages = sum(1 for doc in docs if str((doc.metadata or {}).get("ocr_source") or "") == "paddleocr")
        for doc in docs:
            source = str((doc.metadata or {}).get("ocr_source") or "pypdf")
            doc.metadata = {**(doc.metadata or {}), "parser": source}
        return docs, [
            ParserProvenanceStep(
                name="pypdf+paddleocr",
                status="success" if docs else "empty",
                version=f"pypdf:{_package_version('pypdf')} paddleocr:{_package_version('paddleocr')}",
                stats={"ocr_page_count": ocr_pages, "page_document_count": len(docs)},
            )
        ], ocr_pages

    from langchain_community.document_loaders import PyPDFLoader

    docs = PyPDFLoader(str(path)).load()
    for doc in docs:
        doc.metadata = {**(doc.metadata or {}), "parser": "pypdf"}
    return docs, [ParserProvenanceStep(name="pypdf", status="success", version=_package_version("pypdf"))], 0


def _load_docx(path: Path) -> tuple[list[Document], list[ParserProvenanceStep], int]:
    errors: list[str] = []
    try:
        from langchain_community.document_loaders import Docx2txtLoader

        docs = Docx2txtLoader(str(path)).load()
        if docs:
            for doc in docs:
                doc.metadata = {**(doc.metadata or {}), "parser": "docx2txt"}
            return docs, [ParserProvenanceStep(name="docx2txt", status="success", version=_package_version("docx2txt"))], 0
    except Exception as exc:
        errors.append(f"docx2txt: {exc}")

    try:
        from docx import Document as DocxDocument

        doc = DocxDocument(str(path))
        parts: list[str] = []
        for paragraph in doc.paragraphs:
            text = str(paragraph.text or "").strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [str(cell.text or "").strip() for cell in row.cells if str(cell.text or "").strip()]
                if cells:
                    parts.append(" | ".join(cells))
        content = "\n\n".join(parts).strip()
        if content:
            metadata = {"parser": "python-docx", "document_properties": _core_properties(doc.core_properties)}
            return [Document(page_content=content, metadata=metadata)], [
                ParserProvenanceStep(name="python-docx", status="success", version=_package_version("python-docx"))
            ], 0
    except Exception as exc:
        errors.append(f"python-docx: {exc}")
    raise RuntimeError("; ".join(errors) if errors else "DOCX extraction returned no content.")


def _load_pptx(path: Path) -> tuple[list[Document], list[ParserProvenanceStep], int]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    props = _core_properties(presentation.core_properties)
    docs: list[Document] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = [f"Slide {slide_index}"]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _clean_text(getattr(shape, "text", ""))
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows, start=1):
                    cells = [_clean_text(cell.text) for cell in row.cells if _clean_text(cell.text)]
                    if cells:
                        parts.append(f"Table row {row_index}: " + " | ".join(cells))
        notes = _slide_notes_text(slide)
        if notes:
            parts.append(f"Speaker notes: {notes}")
        content = "\n\n".join(part for part in parts if part.strip()).strip()
        if content:
            docs.append(
                Document(
                    page_content=content,
                    metadata={
                        "parser": "python-pptx",
                        "chunk_index": slide_index - 1,
                        "chunk_type": "slide",
                        "page": slide_index,
                        "slide_number": slide_index,
                        "section_title": f"Slide {slide_index}",
                        "document_properties": props,
                    },
                )
            )
    return docs, [
        ParserProvenanceStep(
            name="python-pptx",
            status="success" if docs else "empty",
            version=_package_version("python-pptx"),
            stats={"slide_count": len(presentation.slides)},
        )
    ], 0


def _load_workbook(path: Path) -> tuple[list[Document], list[ParserProvenanceStep], int]:
    docs = load_workbook_documents(path)
    parser = "pandas" if path.suffix.lower() == ".xls" else "openpyxl"
    for doc in docs:
        doc.metadata = {**(doc.metadata or {}), "parser": parser}
    return docs, [
        ParserProvenanceStep(
            name=parser,
            status="success" if docs else "empty",
            version=_package_version(parser),
            stats={"document_count": len(docs)},
        )
    ], 0


def _load_image(path: Path, settings: Settings) -> tuple[list[Document], list[ParserProvenanceStep], int]:
    if not bool(getattr(settings, "ocr_enabled", False)):
        return [], [
            ParserProvenanceStep(
                name="paddleocr",
                status="skipped",
                version=_package_version("paddleocr"),
                warnings=["OCR disabled; image file was not extracted."],
            )
        ], 0
    docs = load_image_documents(
        path,
        language=str(getattr(settings, "ocr_language", "en") or "en"),
        use_gpu=bool(getattr(settings, "ocr_use_gpu", False)),
    )
    for doc in docs:
        doc.metadata = {**(doc.metadata or {}), "parser": "paddleocr"}
    return docs, [
        ParserProvenanceStep(
            name="paddleocr",
            status="success" if docs else "empty",
            version=_package_version("paddleocr"),
            stats={"image_document_count": len(docs)},
        )
    ], len(docs)


def _core_properties(props: Any) -> dict[str, str]:
    fields = [
        "author",
        "category",
        "comments",
        "content_status",
        "created",
        "identifier",
        "keywords",
        "language",
        "last_modified_by",
        "modified",
        "revision",
        "subject",
        "title",
        "version",
    ]
    payload: dict[str, str] = {}
    for field_name in fields:
        value = getattr(props, field_name, None)
        if value not in ("", None):
            payload[field_name] = str(value)
    return payload


def _slide_notes_text(slide: Any) -> str:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        return _clean_text(getattr(slide.notes_slide.notes_text_frame, "text", "") or "")
    except Exception:
        return ""


def _package_version(package: str) -> str:
    try:
        return importlib.metadata.version(package)
    except Exception:
        return ""


def _clean_text(value: Any) -> str:
    return " ".join(str(value or "").split())


def _unique(items: list[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for item in items:
        text = str(item or "").strip()
        if not text or text in seen:
            continue
        seen.add(text)
        result.append(text)
    return result


def _column_label(index: int) -> str:
    index = int(index)
    letters = ""
    while True:
        index, remainder = divmod(index, 26)
        letters = chr(ord("A") + remainder) + letters
        if index == 0:
            break
        index -= 1
    return letters
