from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace

from agentic_chatbot_next.documents.templates import TemplateTransformService
from agentic_chatbot_next.sandbox.workspace import SessionWorkspace
from agentic_chatbot_next.tools.document_tools import make_document_tools
from agentic_chatbot_next.tools.registry import build_tool_definitions


def _workspace(tmp_path: Path) -> SessionWorkspace:
    workspace = SessionWorkspace.for_session("session-template-transform", tmp_path / "workspaces")
    workspace.open()
    return workspace


def _session(workspace: SessionWorkspace) -> SimpleNamespace:
    return SimpleNamespace(
        tenant_id="tenant",
        user_id="user",
        session_id="session-template-transform",
        conversation_id="conversation-template-transform",
        metadata={},
        workspace=workspace,
    )


def _settings(tmp_path: Path) -> SimpleNamespace:
    return SimpleNamespace(
        default_tenant_id="tenant",
        default_collection_id="default",
        docling_enabled=False,
        workspace_dir=tmp_path / "workspaces",
    )


def _stores() -> SimpleNamespace:
    return SimpleNamespace(doc_store=SimpleNamespace(), chunk_store=SimpleNamespace())


def _write_workspace_doc(workspace: SessionWorkspace, tmp_path: Path, filename: str, body: str) -> None:
    source = tmp_path / filename
    source.write_text(body, encoding="utf-8")
    workspace.copy_file(source)


SOURCE_TEXT = """UNCLASSIFIED

# Program Sustainment Requirements

The contractor shall submit CDRL A001 Program Management Report monthly in PDF format.

The supplier must maintain a verification log for all system tests.

Test setup uses chamber T-5 and calibrated sensor suite S-22.

Test results showed pass for thermal cycle and one anomaly for vibration restart.

The program manager shall review open risks before release.
"""


def _artifact_path(workspace: SessionWorkspace, artifacts: list[dict], suffix: str) -> Path:
    filename = next(item["filename"] for item in artifacts if item["filename"].endswith(suffix))
    return workspace.root / filename


def test_template_transform_memo_creates_docx_markdown_json_and_trace(tmp_path: Path) -> None:
    workspace = _workspace(tmp_path)
    session = _session(workspace)
    _write_workspace_doc(workspace, tmp_path, "sustainment.txt", SOURCE_TEXT)

    result = TemplateTransformService(_settings(tmp_path), _stores(), session).transform(
        document_refs=["sustainment.txt"],
        source_scope="workspace",
        template_type="memo",
        focus="summarize sustainment risks",
        output_format="auto",
    )

    assert result.status == "completed"
    assert result.sections
    assert result.source_documents[0].classification_markings == ["UNCLASSIFIED"]
    generated_names = {item["filename"] for item in result.artifacts}
    trace_names = {item["filename"] for item in result.source_trace_artifacts}
    assert any(name.endswith(".docx") for name in generated_names)
    assert any(name.endswith(".md") for name in generated_names)
    assert any(name.endswith(".json") for name in generated_names)
    assert any(name.endswith("__source_trace.json") for name in trace_names)

    from docx import Document

    document = Document(str(_artifact_path(workspace, result.artifacts, ".docx")))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "Executive Summary" in text
    assert "Source Notes" in text


def test_template_transform_compliance_matrix_creates_xlsx_with_requirement_rows(tmp_path: Path) -> None:
    workspace = _workspace(tmp_path)
    session = _session(workspace)
    _write_workspace_doc(workspace, tmp_path, "rfp.txt", SOURCE_TEXT)

    result = TemplateTransformService(_settings(tmp_path), _stores(), session).transform(
        document_refs=["rfp.txt"],
        source_scope="workspace",
        template_type="rfp_compliance_matrix",
        output_format="auto",
        template_parameters={"owner": "Contracts", "status": "Open"},
    )

    assert result.tables
    rows = result.tables[0].rows
    assert any("contractor shall submit" in row["Requirement Text"].lower() for row in rows)
    assert any(row["Evidence/Source"].startswith("S01.") for row in rows)

    from openpyxl import load_workbook

    workbook = load_workbook(_artifact_path(workspace, result.artifacts, ".xlsx"))
    sheet = workbook["Compliance Matrix"]
    headers = [cell.value for cell in sheet[1]]
    assert headers[:3] == ["Requirement ID", "Source Document", "Clause/Section"]


def test_template_transform_cdrl_sdrl_tracker_extracts_deliverable_rows(tmp_path: Path) -> None:
    workspace = _workspace(tmp_path)
    session = _session(workspace)
    _write_workspace_doc(workspace, tmp_path, "deliverables.txt", SOURCE_TEXT)

    result = TemplateTransformService(_settings(tmp_path), _stores(), session).transform(
        document_refs=["deliverables.txt"],
        source_scope="workspace",
        template_type="cdrl_sdrl_tracker",
        output_format="auto",
    )

    row = result.tables[0].rows[0]
    assert row["Deliverable ID"] == "CDRL-A001"
    assert row["Frequency/Event"].lower() == "monthly"
    assert row["Format"].lower() == "pdf"


def test_template_transform_executive_brief_creates_pptx(tmp_path: Path) -> None:
    workspace = _workspace(tmp_path)
    session = _session(workspace)
    _write_workspace_doc(workspace, tmp_path, "brief-source.txt", SOURCE_TEXT)

    result = TemplateTransformService(_settings(tmp_path), _stores(), session).transform(
        document_refs=["brief-source.txt"],
        source_scope="workspace",
        template_type="executive_brief",
        output_format="auto",
        template_parameters={"title": "Sustainment Brief"},
    )

    assert result.slides
    from pptx import Presentation

    presentation = Presentation(str(_artifact_path(workspace, result.artifacts, ".pptx")))
    titles = [slide.shapes.title.text for slide in presentation.slides if slide.shapes.title]
    assert "Sustainment Brief" in titles
    assert "Key Findings" in titles


def test_template_transform_tool_returns_compact_payload_and_warns_on_unsupported_format(tmp_path: Path) -> None:
    workspace = _workspace(tmp_path)
    session = _session(workspace)
    _write_workspace_doc(workspace, tmp_path, "memo-source.txt", SOURCE_TEXT)
    tools = {
        tool.name: tool
        for tool in make_document_tools(
            _settings(tmp_path),
            _stores(),
            session,
        )
    }

    result = tools["template_transform"].invoke(
        {
            "document_refs": ["memo-source.txt"],
            "source_scope": "workspace",
            "template_type": "memo",
            "output_format": "pptx",
            "run_in_background": False,
        }
    )

    assert result["status"] == "completed"
    assert result["template_type"] == "memo"
    assert "source_documents" not in result
    assert any("not native" in warning for warning in result["warnings"])
    assert any(item["filename"].endswith(".docx") for item in result["generated_artifacts"])


def test_template_transform_registry_metadata() -> None:
    definitions = build_tool_definitions(None)

    definition = definitions["template_transform"]
    assert definition.read_only is True
    assert definition.requires_workspace is True
    assert definition.background_safe is True
    assert definition.concurrency_key == "document_transform"
