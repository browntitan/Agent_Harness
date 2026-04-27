from __future__ import annotations

import builtins
from pathlib import Path
from types import SimpleNamespace

from agentic_chatbot_next.documents.compare import DocumentComparisonService
from agentic_chatbot_next.documents.extractors import DocumentExtractionService
from agentic_chatbot_next.sandbox.workspace import SessionWorkspace
from agentic_chatbot_next.tools.document_tools import make_document_tools
from agentic_chatbot_next.tools.registry import build_tool_definitions


def _workspace(tmp_path: Path) -> SessionWorkspace:
    workspace = SessionWorkspace.for_session("session-docs", tmp_path / "workspaces")
    workspace.open()
    return workspace


def _session(workspace: SessionWorkspace) -> SimpleNamespace:
    return SimpleNamespace(
        tenant_id="tenant",
        session_id="session-docs",
        conversation_id="conversation-docs",
        metadata={},
        workspace=workspace,
    )


def _settings(tmp_path: Path, *, docling_enabled: bool = False) -> SimpleNamespace:
    return SimpleNamespace(
        default_tenant_id="tenant",
        default_collection_id="default",
        docling_enabled=docling_enabled,
        workspace_dir=tmp_path / "workspaces",
    )


def _stores() -> SimpleNamespace:
    return SimpleNamespace(doc_store=SimpleNamespace(), chunk_store=SimpleNamespace())


def _make_docx(path: Path, *, obligation: str = "The contractor shall submit reports monthly.") -> None:
    from docx import Document

    document = Document()
    document.core_properties.title = "Test SOW"
    document.add_heading("Scope", level=1)
    document.add_paragraph(obligation)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Clause"
    table.cell(0, 1).text = "Owner"
    table.cell(1, 0).text = "A.1"
    table.cell(1, 1).text = "Contractor"
    document.save(path)


def _make_pptx(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Program Review"
    textbox = slide.shapes.add_textbox(100000, 900000, 6000000, 600000)
    textbox.text = "The supplier must maintain the delivery log."
    table_shape = slide.shapes.add_table(2, 2, 100000, 1600000, 4000000, 1000000)
    table = table_shape.table
    table.cell(0, 0).text = "Milestone"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "CDR"
    table.cell(1, 1).text = "Open"
    presentation.save(path)


def _make_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "IMS"
    sheet.append(["Milestone", "Current Date"])
    sheet.append(["CDR", "2028-09-26"])
    workbook.save(path)


def _make_pdf(path: Path) -> None:
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    with path.open("wb") as handle:
        writer.write(handle)


def test_document_extract_supports_office_and_pdf_workspace_files(tmp_path: Path) -> None:
    workspace = _workspace(tmp_path)
    session = _session(workspace)
    settings = _settings(tmp_path)
    for filename, maker, expected_parser in [
        ("sow.docx", _make_docx, "python_docx"),
        ("brief.pptx", _make_pptx, "python_pptx"),
        ("tracker.xlsx", _make_xlsx, "openpyxl"),
        ("blank.pdf", _make_pdf, "pypdf_text"),
    ]:
        source = tmp_path / filename
        maker(source)
        workspace.copy_file(source)
        result = DocumentExtractionService(settings, _stores(), session).extract(
            document_ref=filename,
            source_scope="workspace",
            max_elements=50,
        )

        assert result.document.title == filename
        assert result.document.parser_path == expected_parser
        assert result.counts()["metadata_fields"] >= 0
        assert result.elements or result.metadata


def test_document_extract_falls_back_when_docling_enabled_but_missing(tmp_path: Path, monkeypatch) -> None:
    workspace = _workspace(tmp_path)
    session = _session(workspace)
    source = tmp_path / "blank.pdf"
    _make_pdf(source)
    workspace.copy_file(source)

    real_import = builtins.__import__

    def fake_import(name, globals=None, locals=None, fromlist=(), level=0):
        if name.startswith("docling"):
            raise ImportError("docling intentionally absent from slim image")
        return real_import(name, globals, locals, fromlist, level)

    monkeypatch.setattr(builtins, "__import__", fake_import)

    result = DocumentExtractionService(
        _settings(tmp_path, docling_enabled=True),
        _stores(),
        session,
    ).extract(document_ref="blank.pdf", source_scope="workspace", max_elements=50)

    assert result.document.parser_path == "pypdf_text"


def test_document_compare_extracts_changed_obligations_from_workspace_docs(tmp_path: Path) -> None:
    workspace = _workspace(tmp_path)
    session = _session(workspace)
    left = tmp_path / "sow_v1.docx"
    right = tmp_path / "sow_v2.docx"
    _make_docx(left, obligation="The contractor shall submit reports monthly.")
    _make_docx(right, obligation="The contractor may submit reports monthly.")
    workspace.copy_file(left)
    workspace.copy_file(right)

    result = DocumentComparisonService(_settings(tmp_path), _stores(), session).compare(
        left_document_ref="sow_v1.docx",
        right_document_ref="sow_v2.docx",
        source_scope="workspace",
        compare_mode="obligations",
    )

    assert result.counts()["modified"] >= 1
    assert result.changed_obligations
    assert result.changed_obligations[0].change_type == "weakened"
    assert result.changed_obligations[0].severity == "high"


def test_document_extract_auto_scope_falls_back_from_uploads_to_kb(tmp_path: Path) -> None:
    kb_path = tmp_path / "Policy.txt"
    kb_path.write_text("1 Scope\n\nThe contractor shall preserve records.", encoding="utf-8")
    records = {
        "upload-a": SimpleNamespace(
            doc_id="upload-a",
            title="Attachment.txt",
            source_type="upload",
            source_path="",
            collection_id="upload-session",
            file_type="txt",
            content_hash="upload-hash",
            num_chunks=0,
            doc_structure_type="general",
        ),
        "kb-a": SimpleNamespace(
            doc_id="kb-a",
            title="Policy.txt",
            source_type="kb",
            source_path=str(kb_path),
            collection_id="default",
            file_type="txt",
            content_hash="kb-hash",
            num_chunks=0,
            doc_structure_type="general",
        ),
    }
    stores = SimpleNamespace(
        doc_store=SimpleNamespace(
            list_documents=lambda tenant_id="tenant", collection_id="", source_type="": [
                record
                for record in records.values()
                if (not collection_id or record.collection_id == collection_id)
                and (not source_type or record.source_type == source_type)
            ],
            fuzzy_search_title=lambda hint, tenant_id, limit=5, collection_id="": [],
            get_document=lambda doc_id, tenant_id="tenant": records.get(doc_id),
        ),
        chunk_store=SimpleNamespace(list_document_chunks=lambda doc_id, tenant_id="tenant": []),
    )
    session = SimpleNamespace(
        tenant_id="tenant",
        session_id="session-docs",
        conversation_id="conversation-docs",
        metadata={
            "uploaded_doc_ids": ["upload-a"],
            "upload_collection_id": "upload-session",
            "kb_collection_id": "default",
        },
    )

    result = DocumentExtractionService(_settings(tmp_path), stores, session).extract(
        document_ref="Policy.txt",
        source_scope="auto",
        max_elements=20,
    )

    assert result.document.doc_id == "kb-a"
    assert result.document.source_scope == "kb"
    assert any("preserve records" in element.text for element in result.elements)


def test_document_tools_return_compact_payloads_and_artifacts(tmp_path: Path) -> None:
    workspace = _workspace(tmp_path)
    session = _session(workspace)
    source = tmp_path / "sow.docx"
    _make_docx(source)
    workspace.copy_file(source)

    tools = {
        tool.name: tool
        for tool in make_document_tools(
            _settings(tmp_path),
            _stores(),
            session,
        )
    }

    result = tools["document_extract"].invoke(
        {
            "document_ref": "sow.docx",
            "source_scope": "workspace",
            "output": "summary",
            "export": True,
        }
    )

    assert result["counts"]["sections"] >= 1
    assert "elements" not in result
    assert len(result["artifacts"]) == 2
    assert all(item["artifact_ref"].startswith("download://") for item in result["artifacts"])


def test_document_registry_metadata_marks_tools_read_only_workspace_and_background_safe() -> None:
    definitions = build_tool_definitions(None)

    assert definitions["document_extract"].read_only is True
    assert definitions["document_extract"].requires_workspace is True
    assert definitions["document_extract"].background_safe is True
    assert definitions["document_extract"].concurrency_key == "document_analysis"
    assert definitions["document_compare"].read_only is True
    assert definitions["document_compare"].requires_workspace is True
    assert definitions["document_compare"].background_safe is True
    assert definitions["document_compare"].concurrency_key == "document_analysis"
    assert definitions["document_consolidation_campaign"].read_only is True
    assert definitions["document_consolidation_campaign"].requires_workspace is True
    assert definitions["document_consolidation_campaign"].background_safe is True
    assert definitions["document_consolidation_campaign"].concurrency_key == "document_campaign"
    assert definitions["template_transform"].read_only is True
    assert definitions["template_transform"].requires_workspace is True
    assert definitions["template_transform"].background_safe is True
    assert definitions["template_transform"].concurrency_key == "document_transform"
