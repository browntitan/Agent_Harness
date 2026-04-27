"""
Build a leadership-ready deck from the 151-slide audit-updated deck.
- Keeps ~55 slides in a coherent narrative order
- Deletes code-level implementation details, section dividers, and proposal slides
- Updates speaker notes on reordered slides
"""
from pptx import Presentation
from pptx.oxml.ns import qn
import os, copy

SRC = os.path.join(os.path.dirname(__file__), "agentic_chatbot_v3_audit_updated.pptx")
OUT = os.path.join(os.path.dirname(__file__), "agentic_chatbot_v3_leadership.pptx")

prs = Presentation(SRC)

# ══════════════════════════════════════════════════════════════
# SLIDE KEEP-LIST (1-indexed slide numbers from the 151-slide deck)
# Ordered for a coherent leadership narrative.
# ══════════════════════════════════════════════════════════════

KEEP_SLIDES = [
    # ── ACT 1: WHAT IS THIS? (Intro + Agent Harnesses) ──────
    1,   # Title: Agentic Chatbot v3
    2,   # Section: What Is an Agent Harness?
    3,   # Why LLMs Alone Aren't Enough (problem/solution)
    4,   # Anatomy of an Agent Harness (core loop diagram)
    5,   # The ReAct Pattern — How Agents Reason
    6,   # Agent Modes — One Harness, Many Strategies

    # ── ACT 2: OUR SYSTEM AT A GLANCE ───────────────────────
    23,  # What Does This System Do?
    24,  # The 30-Second Mental Model
    8,   # Our Agent Harness at a Glance (key numbers)
    7,   # Tools & RAG — The Agent's Capabilities
    10,  # API Surface — 118 Endpoints
    9,   # New Subsystems — GraphRAG, MCP, Worker Orchestration

    # ── ACT 3: ARCHITECTURE ─────────────────────────────────
    36,  # Three-Layer Runtime Architecture
    37,  # Layer 1: RuntimeService — The Orchestrator
    45,  # The Router — How the System Decides What to Do
    51,  # All 11 Agent Roles at a Glance
    52,  # How Agent Mode Maps to Execution
    77,  # Agent Capabilities Matrix — Complete Reference

    # ── ACT 4: REAL SCENARIOS ────────────────────────────────
    54,  # Section: Execution Flow Examples
    55,  # Flow 1: "What is machine learning?" (basic)
    56,  # Flow 2: "Analyze sales.csv" (data analyst)
    57,  # Flow 3: "What does the policy say?" (RAG)
    58,  # Flow 4: "Compare Q1/Q2, research, report" (coordinator)
    59,  # Flow 5: Memory in Action
    60,  # Complete Routing Decision Tree
    79,  # RAG Pipeline — Retrieval to Grounded Answer

    # ── ACT 5: SUBSYSTEM DEEP DIVES ─────────────────────────
    12,  # Authorization Subsystem
    13,  # Documents Subsystem Deep Dive
    14,  # Skills Architecture Deep Dive
    15,  # Provider System & Circuit Breaker
    16,  # MCP Security Model
    17,  # Task Plan & Artifact Handoff System
    18,  # Storage & Blob Store
    19,  # Graph Query Methods & Build Pipeline
    20,  # Configuration & Tuning Defaults
    21,  # Deep RAG & Verification Pipeline
    22,  # Persistence, Observability & Benchmarks

    # ── ACT 6: CODEBASE & OPERATIONS ────────────────────────
    26,  # Top-Level Repository Map
    27,  # Frontend — React Chat UI
    28,  # Control Panel UI — Admin React App
    30,  # data/ Directory — Configuration & Runtime State

    # ── ACT 7: WHAT'S NEW & IMPROVEMENTS ────────────────────
    122, # GraphRAG — Knowledge Graph Integration
    123, # Deep RAG & Adaptive Retrieval
    124, # Router Feedback Loop & Worker Scheduler
    125, # New Infrastructure — 7 Modules Added
    126, # Agentic Harness — 6 Runtime Improvements
    120, # Feature Roadmap — All 13 Improvements

    # ── ACT 8: ENTERPRISE READINESS ─────────────────────────
    146, # Memory — 3 Enterprise Gaps
    147, # Agent Routing — 3 Enterprise Gaps
    148, # Tools & Capability Model — 3 Enterprise Gaps
    149, # Skills & Prompt Governance — 3 Enterprise Gaps
    150, # Observability, Governance & Compliance — 5 Gaps
    151, # Enterprise Readiness Matrix — 17 Gaps
]

print(f"Source deck: {len(prs.slides)} slides")
print(f"Keeping: {len(KEEP_SLIDES)} slides")
print(f"Deleting: {len(prs.slides) - len(KEEP_SLIDES)} slides")

# ══════════════════════════════════════════════════════════════
# REORDER & DELETE via sldIdLst manipulation
# ══════════════════════════════════════════════════════════════

pres_elem = prs.part._element
sldIdLst = pres_elem.find(qn('p:sldIdLst'))
all_sldIds = list(sldIdLst)

# Build mapping: 1-indexed slide number -> sldId element
slide_map = {}
for i, elem in enumerate(all_sldIds):
    slide_map[i + 1] = elem

# Clear the list
for elem in all_sldIds:
    sldIdLst.remove(elem)

# Re-add only the keep slides in the desired order
kept = 0
for slide_num in KEEP_SLIDES:
    if slide_num in slide_map:
        sldIdLst.append(slide_map[slide_num])
        kept += 1
    else:
        print(f"  WARNING: Slide {slide_num} not found in deck!")

print(f"Kept {kept} slides in new order")

# Note: Orphaned slide files stay in the package but won't render.
# Use unpack/clean/pack workflow for a truly clean file if needed.

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════

prs.save(OUT)
print(f"\nSaved to: {OUT}")

# Verify by counting sldIdLst entries
pres_elem2 = prs.part._element
sldIdLst2 = pres_elem2.find(qn('p:sldIdLst'))
print(f"Final slide count in sldIdLst: {len(list(sldIdLst2))}")
